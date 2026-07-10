from __future__ import annotations

import ast
import base64
import binascii
import hashlib
import html
import ipaddress
import json
import mimetypes
import re
import socket
import time
from collections import defaultdict
from pathlib import Path
from typing import Any, Callable
from urllib.parse import parse_qs, urlsplit

import logging

import httpx
from sqlalchemy.orm import Session

from app.core.config import get_settings
from app.core.enums import MaterialType
from app.core.errors import bad_request
from app.core.media import normalize_storage_path, signed_media_url
from app.services.runtime_config import get_enabled_service_config
from app.services.storage import storage_service


LOGGER = logging.getLogger(__name__)


SERVICE_TYPE = "doc_parser"
DEFAULT_ENDPOINT = "docmind-api.cn-hangzhou.aliyuncs.com"
DEFAULT_REGION = "cn-hangzhou"
DEFAULT_LAYOUT_STEP_SIZE = 100
DEFAULT_DOC_PARSER_TIMEOUT_SECONDS = 7200
MAX_DOC_PARSER_TIMEOUT_SECONDS = 7200
DEFAULT_DOC_PARSER_POLL_INTERVAL_SECONDS = 10
SUPPORTED_MATERIAL_TYPES = {MaterialType.PPTX.value, MaterialType.PDF.value, MaterialType.DOCX.value, MaterialType.TXT.value, MaterialType.IMAGE.value}
MARKDOWN_IMAGE_URL_PATTERN = re.compile(
    r"!\[(?P<alt>[^\]]*)]\(\s*(?:<(?P<angle>[^>]+)>|(?P<plain>[^)\s]+))(?:\s+['\"][^'\"]*['\"])?\s*\)"
)
HTML_IMAGE_URL_PATTERN = re.compile(r"""<img\b[^>]*\bsrc=["'](?P<src>[^"']+)["'][^>]*>""", re.IGNORECASE)
DATA_IMAGE_PATTERN = re.compile(r"^data:(?P<mime>image/[a-z0-9.+-]+);base64,(?P<data>.+)$", re.IGNORECASE | re.DOTALL)
MAX_MARKDOWN_IMAGE_BYTES = 15 * 1024 * 1024
MAX_MARKDOWN_RESULT_BYTES = 20 * 1024 * 1024
SIGNED_IMAGE_QUERY_KEYS = {"Expires", "OSSAccessKeyId", "Signature", "security-token", "x-oss-security-token"}
IMAGE_FILENAME_ALT_PATTERN = re.compile(r"^[A-Fa-f0-9_-]{12,}\.(?:png|jpe?g|gif|webp|bmp|svg)$", re.IGNORECASE)
MAX_MARKDOWN_IMAGE_REDIRECTS = 3


def _normalize_page(title: str | None, content: str, page_number: int) -> dict:
    clean_content = content.strip()
    return {
        "page_number": page_number,
        "page_title": title.strip() if title else None,
        "page_text": clean_content or "本页未提取到有效文字内容。",
    }


def _to_plain_data(value: Any) -> Any:
    if value is None:
        return None
    if hasattr(value, "to_map"):
        return value.to_map()
    if isinstance(value, str):
        text = value.strip()
        if not text:
            return {}
        if text.startswith("{") or text.startswith("["):
            try:
                return json.loads(text)
            except json.JSONDecodeError:
                return {"text": text}
        return {"text": text}
    if isinstance(value, list):
        return [_to_plain_data(item) for item in value]
    if isinstance(value, dict):
        return {key: _to_plain_data(item) for key, item in value.items()}
    return value


def _read_value(value: Any, *keys: str) -> Any:
    if value is None:
        return None
    for key in keys:
        if isinstance(value, dict):
            for candidate in {key, key[:1].upper() + key[1:], key[:1].lower() + key[1:]}:
                if candidate in value:
                    return value[candidate]
        if hasattr(value, key):
            return getattr(value, key)
    return None


def _as_bool(value: Any, default: bool = False) -> bool:
    if value is None:
        return default
    if isinstance(value, bool):
        return value
    if isinstance(value, (int, float)):
        return bool(value)
    return str(value).strip().lower() in {"1", "true", "yes", "on", "启用", "是"}


def _bounded_int(value: Any, *, default: int, minimum: int, maximum: int) -> int:
    try:
        number = int(value)
    except (TypeError, ValueError):
        number = default
    return max(minimum, min(maximum, number))


def _safe_int(value: Any, default: int = 0) -> int:
    try:
        return int(value)
    except (TypeError, ValueError):
        return default


def _api_error_message(body: Any) -> str | None:
    data = _to_plain_data(body)
    if not isinstance(data, dict):
        return None
    code = data.get("Code") or data.get("code")
    if code is None:
        return None
    if str(code).lower() in {"200", "ok", "success"}:
        return None
    return str(data.get("Message") or data.get("message") or code)


def _output_formats(value: Any) -> list[str] | None:
    if value is None or value == "":
        return ["markdown"]
    if isinstance(value, str):
        items = [item.strip() for item in value.split(",")]
    elif isinstance(value, list):
        items = [str(item).strip() for item in value]
    else:
        items = [str(value).strip()]
    formats = [item for item in items if item]
    return formats or None


TEXT_PAYLOAD_KEYS = ("markdownContent", "markdown_content", "llmResult", "llm_result", "page_text", "script_text", "content", "text")


def _decode_serialized_payload(value: Any) -> Any:
    if not isinstance(value, str):
        return value
    text = value.strip()
    if not text or text[0] not in "{[":
        return text
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        try:
            decoded = ast.literal_eval(text)
        except (SyntaxError, ValueError):
            return text
        return decoded if isinstance(decoded, (dict, list, str)) else text


def _extract_serialized_text_values(value: str) -> str:
    key_pattern = re.compile(r"['\"](?:" + "|".join(TEXT_PAYLOAD_KEYS) + r")['\"]\s*:\s*")
    pieces: list[str] = []
    cursor = 0
    while True:
        match = key_pattern.search(value, cursor)
        if match is None:
            break
        cursor = match.end()
        while cursor < len(value) and value[cursor].isspace():
            cursor += 1
        if cursor >= len(value) or value[cursor] not in {"'", '"'}:
            continue
        quote = value[cursor]
        cursor += 1
        raw = ""
        escaped = False
        while cursor < len(value):
            char = value[cursor]
            cursor += 1
            if escaped:
                raw += f"\\{char}"
                escaped = False
                continue
            if char == "\\":
                escaped = True
                continue
            if char == quote:
                if raw.strip():
                    pieces.append(raw)
                break
            raw += char
        if escaped and raw.strip():
            pieces.append(f"{raw}\\")
    return "\n\n".join(pieces)


def _extract_text_payload(value: Any) -> str:
    value = _decode_serialized_payload(value)
    if value is None:
        return ""
    if hasattr(value, "to_map"):
        value = value.to_map()
    if isinstance(value, str):
        extracted = _extract_serialized_text_values(value) if value.lstrip().startswith(("{", "[")) else ""
        text = extracted or value
        return text.replace("\\n", "\n").replace("\\t", "\t").replace("\\'", "'").replace('\\"', '"').strip()
    if isinstance(value, list):
        pieces = [_extract_text_payload(item) for item in value]
        return "\n\n".join(piece for piece in pieces if piece)
    if isinstance(value, dict):
        blocks = value.get("blocks")
        if isinstance(blocks, list):
            pieces = [_extract_text_payload(block) for block in blocks]
            text = "\n".join(piece for piece in pieces if piece)
            if text:
                return text
        for key in TEXT_PAYLOAD_KEYS:
            if key in value:
                text = _extract_text_payload(value.get(key))
                if text:
                    return text
        pieces = [_extract_text_payload(item) for item in value.values()]
        return "\n\n".join(piece for piece in pieces if piece)
    return str(value).strip()


def _layout_text(layout: dict) -> str:
    for key in TEXT_PAYLOAD_KEYS:
        if key in layout:
            text = _extract_text_payload(layout.get(key))
            if text:
                return text
    return _extract_text_payload(layout.get("blocks"))


def _layout_title(layout: dict, fallback: str | None = None) -> str | None:
    text = str(layout.get("text") or fallback or "").strip()
    if not text:
        content = _layout_text(layout)
        text = content.splitlines()[0].strip("# ").strip() if content else ""
    return text[:120] or None


def _markdown_title(content: str, fallback: str = "文档内容") -> str:
    for line in content.splitlines():
        clean = line.strip()
        if clean.startswith("#"):
            title = clean.lstrip("#").strip()
            if title:
                return title[:120]
        if clean:
            return clean[:120]
    return fallback


def _pages_from_markdown(content: str) -> list[dict]:
    text = content.strip()
    if not text:
        return []
    page_markers = list(re.finditer(r"(?m)^\s*(?:<!--\s*pagebreak\s*-->|-{3,}\s*$)", text))
    if not page_markers:
        return [_normalize_page(_markdown_title(text), text, 1)]
    pages: list[dict] = []
    cursor = 0
    for marker in page_markers:
        chunk = text[cursor : marker.start()].strip()
        if chunk:
            pages.append(_normalize_page(_markdown_title(chunk, f"第{len(pages) + 1}页"), chunk, len(pages) + 1))
        cursor = marker.end()
    tail = text[cursor:].strip()
    if tail:
        pages.append(_normalize_page(_markdown_title(tail, f"第{len(pages) + 1}页"), tail, len(pages) + 1))
    return pages


def _parse_image(path: Path, db: Session | None) -> list[dict]:
    # 修复 DEF-02：图片课程资料经 OCR 识别为文本，生成单页课时内容；OCR 未配置或未识别到文本时
    # 抛出可理解的错误，由上层将资料标记为 failed（诚实降级，不产出空课时）。
    from app.services.ocr import ocr_service

    content = path.read_bytes()
    text = (ocr_service.recognize_bytes(content, db=db) or "").strip()
    if not text:
        raise bad_request("图片资料未识别到文本，请上传更清晰的图片")
    return [_normalize_page(_markdown_title(text, "图片资料"), text, 1)]


def _parse_local_fallback(path: Path, material_type: str) -> list[dict]:
    if material_type == MaterialType.TXT.value:
        text = path.read_text(encoding="utf-8", errors="ignore")
        chunks = [chunk.strip() for chunk in re.split(r"\n\s*\n+", text) if chunk.strip()]
        if len(chunks) > 1:
            return [_normalize_page(_markdown_title(chunk, f"第{index}页"), chunk, index) for index, chunk in enumerate(chunks, start=1)]
        return _pages_from_markdown(text)
    if material_type == MaterialType.PPTX.value:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise bad_request("缺少 python-pptx 依赖，无法本地解析 PPTX") from exc
        presentation = Presentation(path)
        pages: list[dict] = []
        for index, slide in enumerate(presentation.slides, start=1):
            pieces = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            pages.append(_normalize_page(pieces[0] if pieces else f"第{index}页", "\n\n".join(pieces), index))
        return pages
    if material_type == MaterialType.DOCX.value:
        try:
            from docx import Document
        except ImportError as exc:
            raise bad_request("缺少 python-docx 依赖，无法本地解析 DOCX") from exc
        document = Document(path)
        text = "\n\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())
        return _pages_from_markdown(text)
    if material_type == MaterialType.PDF.value:
        try:
            import fitz
        except ImportError as exc:
            raise bad_request("缺少 pymupdf 依赖，无法本地解析 PDF") from exc
        pages = []
        with fitz.open(path) as document:
            for index, page in enumerate(document, start=1):
                text = page.get_text("text").strip()
                pages.append(_normalize_page(f"第{index}页", text, index))
        return pages
    return []


def _markdown_image_urls(content: str) -> list[str]:
    urls: list[str] = []
    seen: set[str] = set()
    for pattern in (MARKDOWN_IMAGE_URL_PATTERN, HTML_IMAGE_URL_PATTERN):
        for match in pattern.finditer(content):
            raw = match.groupdict().get("angle") or match.groupdict().get("plain") or match.groupdict().get("src")
            if not raw:
                continue
            value = raw.strip()
            if value and value not in seen:
                urls.append(value)
                seen.add(value)
    return urls


def _is_remote_url(value: str) -> bool:
    try:
        parsed = urlsplit(html.unescape(value))
    except ValueError:
        return False
    return parsed.scheme in {"http", "https"} and bool(parsed.netloc) and _url_host_allowed(value)


def _ip_allowed(ip: ipaddress._BaseAddress) -> bool:
    # 用正向 is_global 判定（仅放行全球可路由地址），自动排除私网/环回/链路本地/保留/多播/未指定，
    # 并覆盖 IPv4-mapped IPv6(如 ::ffff:127.0.0.1)等否定式黑名单易漏的段。
    mapped = getattr(ip, "ipv4_mapped", None)
    if mapped is not None:
        ip = mapped
    return bool(getattr(ip, "is_global", False)) and not (
        ip.is_private or ip.is_loopback or ip.is_link_local or ip.is_multicast or ip.is_reserved or ip.is_unspecified
    )


def _url_host_allowed(value: str) -> bool:
    try:
        parsed = urlsplit(html.unescape(value))
    except ValueError:
        return False
    hostname = parsed.hostname
    if not hostname:
        return False
    try:
        direct_ip = ipaddress.ip_address(hostname)
    except ValueError:
        direct_ip = None
    if direct_ip is not None:
        return _ip_allowed(direct_ip)
    try:
        infos = socket.getaddrinfo(hostname, parsed.port or (443 if parsed.scheme == "https" else 80), type=socket.SOCK_STREAM)
    except socket.gaierror:
        return False
    addresses = {item[4][0] for item in infos}
    if not addresses:
        return False
    for address in addresses:
        try:
            if not _ip_allowed(ipaddress.ip_address(address)):
                return False
        except ValueError:
            return False
    return True


def _resolve_allowed_ip(hostname: str, port: int, scheme: str) -> str | None:
    # M14: 解析 host 得到 IP，并对每个解析结果做私网/保留/元数据网段拒绝校验。
    # 返回首个校验通过的 IP（literal IP 直接校验返回）。任一解析结果不通过则整体拒绝，
    # 避免攻击者在 DNS 中混入内网地址。
    try:
        direct_ip = ipaddress.ip_address(hostname)
    except ValueError:
        direct_ip = None
    if direct_ip is not None:
        return str(direct_ip) if _ip_allowed(direct_ip) else None
    try:
        infos = socket.getaddrinfo(hostname, port, type=socket.SOCK_STREAM)
    except socket.gaierror:
        return None
    addresses = [item[4][0] for item in infos]
    if not addresses:
        return None
    chosen: str | None = None
    for address in addresses:
        try:
            candidate = ipaddress.ip_address(address)
        except ValueError:
            return None
        if not _ip_allowed(candidate):
            return None
        if chosen is None:
            chosen = str(candidate)
    return chosen


class _PinnedIPTransport(httpx.HTTPTransport):
    """M14 SSRF / DNS-rebinding TOCTOU 防护。

    方案：自定义 transport 拦截每一次请求（含每一跳重定向），在即将连接之前
    自行 socket.getaddrinfo 解析目标 host，复用 _ip_allowed 对解析出的 IP 做
    私网/环回/链路本地/保留/元数据(169.254.169.254) 网段拒绝校验，然后把请求 URL 的
    host 替换为这个“校验通过的 IP”，使 httpcore 直连该 IP——校验与连接锁定为同一个 IP，
    杜绝“校验后再次 DNS 解析”的时间窗。同时：
      - 保留原始 Host 头（Host: 原 hostname），让目标按虚拟主机正确路由；
      - https 通过 sni_hostname 扩展把 TLS SNI / 证书校验的 server_hostname 设为原 hostname
        （而非 IP），保证证书校验仍针对真实域名、verify 不被削弱。
    残留风险：仅在“解析→校验→连接”的极短同步窗口内复用本进程一次解析结果；httpcore 不会
    再次解析（连接目标已是字面 IP），故 rebinding 无法在校验与连接之间切换到内网。
    """

    def handle_request(self, request: httpx.Request) -> httpx.Response:
        url = request.url
        scheme = url.scheme
        if scheme not in {"http", "https"}:
            raise httpx.ConnectError("blocked non-http(s) url", request=request)
        hostname = url.host
        if not hostname:
            raise httpx.ConnectError("missing host", request=request)
        # 若 URL host 已是字面 IP，直接校验该 IP；否则解析域名得到 IP 再校验。
        port = url.port or (443 if scheme == "https" else 80)
        pinned_ip = _resolve_allowed_ip(hostname, port, scheme)
        if pinned_ip is None:
            raise httpx.ConnectError("blocked private/reserved address", request=request)
        try:
            ip_literal = ipaddress.ip_address(pinned_ip)
        except ValueError:
            raise httpx.ConnectError("invalid resolved address", request=request)
        host_is_literal_ip = False
        try:
            ipaddress.ip_address(hostname)
            host_is_literal_ip = True
        except ValueError:
            host_is_literal_ip = False
        # 把连接目标锁定为校验通过的 IP（IPv6 需用 [..] 包裹）。
        connect_host = f"[{pinned_ip}]" if ip_literal.version == 6 else pinned_ip
        request.url = url.copy_with(host=connect_host)
        if not host_is_literal_ip:
            # 保留原始 Host 头（host[:port]，仅在显式非默认端口时带端口），供目标按域名虚拟主机路由。
            default_port = 443 if scheme == "https" else 80
            host_header = hostname if url.port in (None, default_port) else f"{hostname}:{url.port}"
            request.headers["Host"] = host_header
            # https 下让 TLS SNI / 证书校验仍针对真实域名而非字面 IP。
            if scheme == "https":
                request.extensions = {**request.extensions, "sni_hostname": hostname}
        return super().handle_request(request)


def _build_pinned_client(timeout: httpx.Timeout) -> httpx.Client:
    return httpx.Client(
        transport=_PinnedIPTransport(),
        timeout=timeout,
        follow_redirects=False,
        headers={"User-Agent": "class-agent-doc-parser/1.0"},
    )


def _is_temporary_docmind_url(value: str) -> bool:
    try:
        parsed = urlsplit(html.unescape(value))
    except ValueError:
        return False
    if parsed.scheme not in {"http", "https"}:
        return False
    query = parse_qs(parsed.query)
    return bool(SIGNED_IMAGE_QUERY_KEYS.intersection(query))


def _local_docmind_image_url(value: str) -> str | None:
    # H2: 该函数被 Pydantic 序列化器 sanitize_temporary_docmind_images 调用，
    # 绝不能向上抛异常（否则 GET /materials 整页 500，持久 DoS），且绝不能用未规范化、
    # 含 `..`/绝对路径的相对路径去触碰文件系统或参与签名（路径穿越 + signed_media_url 抛 ValueError）。
    try:
        parsed = urlsplit(html.unescape(value))
        if parsed.netloc and not storage_service._is_loopback_url(value):
            return None
        path = parsed.path.lstrip("/")
        marker = "docmind_images/"
        if marker not in path:
            return None
        relative_path = path[path.index(marker) :]
        if not relative_path:
            return None
        # 先用 normalize_storage_path 做规范化校验：拒绝任何含 `..`/绝对路径/越界的目标。
        # urlsplit().path 不会折叠 `..`，必须在触碰文件系统/签名之前拦截，避免跨目录文件
        # 存在性探测（OS stat 真实解析 `..`）以及后续签名时 `..` in path.parts 抛 ValueError。
        try:
            safe_path = normalize_storage_path(relative_path)
        except ValueError:
            return None
        if not safe_path.startswith(marker):
            return None
        if (
            not storage_service.absolute_path(safe_path).is_file()
            and not storage_service.absolute_path(f"public/{safe_path}").is_file()
        ):
            return None
        return signed_media_url(safe_path)
    except Exception:
        # 兜底：任何意外异常都返回 None，确保序列化器永不因本函数 500。
        return None


def _image_suffix(url: str, content_type: str | None) -> str:
    suffix = Path(urlsplit(url).path).suffix.lower()
    if suffix in {".png", ".jpg", ".jpeg", ".webp"}:
        return suffix
    mime = str(content_type or "").split(";", 1)[0].strip().lower()
    guessed = mimetypes.guess_extension(mime) if mime else None
    if guessed == ".jpe":
        return ".jpg"
    return guessed if guessed in {".png", ".jpg", ".jpeg", ".webp"} else ".png"


def _image_content_matches(content: bytes, suffix: str, mime: str | None = None) -> bool:
    normalized_suffix = ".jpg" if suffix == ".jpeg" else suffix
    normalized_mime = (mime or "").split(";", 1)[0].strip().lower()
    if normalized_suffix == ".png" or normalized_mime == "image/png":
        return content.startswith(b"\x89PNG\r\n\x1a\n")
    if normalized_suffix == ".jpg" or normalized_mime == "image/jpeg":
        return content.startswith(b"\xff\xd8\xff")
    if normalized_suffix == ".webp" or normalized_mime == "image/webp":
        return len(content) >= 12 and content.startswith(b"RIFF") and content[8:12] == b"WEBP"
    return False


def _persist_data_image(url: str, *, db: Session | None, cache: dict[str, str | None]) -> str | None:
    data_url = html.unescape(url).strip()
    if data_url in cache:
        return cache[data_url]
    cache[data_url] = None
    match = DATA_IMAGE_PATTERN.match(data_url)
    if match is None:
        return None
    try:
        content = base64.b64decode(match.group("data"), validate=True)
    except (binascii.Error, ValueError):
        return None
    if not content or len(content) > MAX_MARKDOWN_IMAGE_BYTES:
        return None
    digest = hashlib.sha256(content).hexdigest()
    suffix = _image_suffix(f"image.{match.group('mime').split('/', 1)[1]}", match.group("mime"))
    if not _image_content_matches(content, suffix, match.group("mime")):
        return None
    relative_path = storage_service.save_bytes(
        content,
        folder=f"docmind_images/{digest[:2]}",
        filename=f"{digest}{suffix}",
        db=db,
    )
    media_url = signed_media_url(relative_path)
    cache[data_url] = media_url
    return media_url


def _persist_markdown_image(
    client: httpx.Client,
    url: str,
    *,
    db: Session | None,
    cache: dict[str, str | None],
) -> str | None:
    download_url = html.unescape(url)
    if DATA_IMAGE_PATTERN.match(download_url.strip()):
        return _persist_data_image(download_url, db=db, cache=cache)
    if download_url in cache:
        return cache[download_url]
    cache[download_url] = None
    if not _is_remote_url(download_url):
        return None
    current_url = download_url
    response: httpx.Response | None = None
    try:
        for _ in range(MAX_MARKDOWN_IMAGE_REDIRECTS + 1):
            if not _is_remote_url(current_url):
                return None
            response = client.get(current_url)
            if bool(getattr(response, "is_redirect", False)):
                location = response.headers.get("location")
                if not location:
                    return None
                # 用调用方原始 URL（非被 transport 改写成 IP 的 response.url）解析相对重定向，
                # 保证下一跳仍按真实域名重新解析+校验（_PinnedIPTransport 会再次锁 IP）。
                current_url = str(httpx.URL(current_url).join(location))
                continue
            response.raise_for_status()
            break
        else:
            return None
    except httpx.HTTPError:
        return None
    if response is None:
        return None
    content_type = response.headers.get("content-type", "")
    content = b""
    if hasattr(response, "iter_bytes"):
        for chunk in response.iter_bytes(chunk_size=64 * 1024):
            content += chunk
            if len(content) > MAX_MARKDOWN_IMAGE_BYTES:
                return None
    else:
        content = getattr(response, "content", b"")
    if not content or len(content) > MAX_MARKDOWN_IMAGE_BYTES:
        return None
    mime = content_type.split(";", 1)[0].strip().lower()
    if not mime.startswith("image/"):
        suffix = Path(urlsplit(current_url).path).suffix.lower()
        if mime and mime not in {"application/octet-stream", "binary/octet-stream"}:
            return None
        if suffix not in {".png", ".jpg", ".jpeg", ".webp"}:
            return None
    elif mime not in {"image/png", "image/jpeg", "image/webp"}:
        return None
    digest = hashlib.sha256(download_url.encode("utf-8")).hexdigest()
    suffix = _image_suffix(current_url, content_type)
    if not _image_content_matches(content, suffix, content_type):
        return None
    relative_path = storage_service.save_bytes(
        content,
        folder=f"docmind_images/{digest[:2]}",
        filename=f"{digest}{suffix}",
        db=db,
    )
    media_url = signed_media_url(relative_path)
    cache[download_url] = media_url
    return media_url


def _same_image_url(left: str | None, right: str) -> bool:
    if not left:
        return False
    return html.unescape(left).strip() == html.unescape(right).strip()


def _strip_unavailable_image(content: str, raw_url: str) -> str:
    placeholder = "（图片未能保存，请重新解析资料）"

    def display_alt(alt: str) -> str:
        clean = html.unescape(alt).strip()
        if IMAGE_FILENAME_ALT_PATTERN.fullmatch(clean):
            return ""
        return clean

    def replace_markdown(match: re.Match[str]) -> str:
        url = match.groupdict().get("angle") or match.groupdict().get("plain")
        if not _same_image_url(url, raw_url):
            return match.group(0)
        alt = display_alt(match.groupdict().get("alt") or "")
        return alt or placeholder

    def replace_html(match: re.Match[str]) -> str:
        if not _same_image_url(match.groupdict().get("src"), raw_url):
            return match.group(0)
        alt_match = re.search(r"""\balt=["'](?P<alt>[^"']+)["']""", match.group(0), flags=re.IGNORECASE)
        alt = display_alt(alt_match.group("alt") if alt_match else "")
        return alt or placeholder

    return HTML_IMAGE_URL_PATTERN.sub(replace_html, MARKDOWN_IMAGE_URL_PATTERN.sub(replace_markdown, content))


def _replace_image_url(content: str, raw_url: str, new_url: str) -> str:
    def display_alt(alt: str) -> str:
        clean = html.unescape(alt).strip()
        return "课件图片" if IMAGE_FILENAME_ALT_PATTERN.fullmatch(clean) else clean

    def replace_markdown(match: re.Match[str]) -> str:
        url = match.groupdict().get("angle") or match.groupdict().get("plain")
        if not _same_image_url(url, raw_url):
            return match.group(0)
        alt = display_alt(match.groupdict().get("alt") or "")
        return f"![{alt}]({new_url})"

    def replace_html(match: re.Match[str]) -> str:
        if not _same_image_url(match.groupdict().get("src"), raw_url):
            return match.group(0)
        updated = match.group(0).replace(match.groupdict().get("src") or "", new_url)
        alt_match = re.search(r"""\balt=["'](?P<alt>[^"']*)["']""", updated, flags=re.IGNORECASE)
        if alt_match and IMAGE_FILENAME_ALT_PATTERN.fullmatch(html.unescape(alt_match.group("alt")).strip()):
            updated = updated[: alt_match.start("alt")] + "课件图片" + updated[alt_match.end("alt") :]
        return updated

    return HTML_IMAGE_URL_PATTERN.sub(replace_html, MARKDOWN_IMAGE_URL_PATTERN.sub(replace_markdown, content))


def _localize_markdown_images(content: str, db: Session | None, cache: dict[str, str | None]) -> str:
    urls = _markdown_image_urls(content)
    if not urls:
        return content
    rewritten = content
    with _build_pinned_client(httpx.Timeout(20.0, connect=8.0)) as client:
        for raw_url in urls:
            public_url = _persist_markdown_image(client, raw_url, db=db, cache=cache)
            if not public_url:
                if _is_temporary_docmind_url(raw_url):
                    rewritten = _strip_unavailable_image(rewritten, raw_url)
                continue
            rewritten = _replace_image_url(rewritten, raw_url, public_url)
    return rewritten


def sanitize_temporary_docmind_images(content: str | None) -> str | None:
    if not content:
        return content
    rewritten = str(content)
    for raw_url in _markdown_image_urls(rewritten):
        if _is_temporary_docmind_url(raw_url):
            rewritten = _strip_unavailable_image(rewritten, raw_url)
            continue
        local_url = _local_docmind_image_url(raw_url)
        if local_url and local_url != raw_url:
            rewritten = _replace_image_url(rewritten, raw_url, local_url)
        elif "docmind_images/" in html.unescape(raw_url):
            rewritten = _replace_image_url(rewritten, raw_url, raw_url)
    return rewritten


class DocParserService:
    def __init__(self) -> None:
        self.settings = get_settings()

    def _client(self, config: dict):
        required = ["access_key_id", "access_key_secret"]
        missing = [key for key in required if not config.get(key)]
        if missing:
            raise bad_request(f"文档解析配置缺少字段: {', '.join(missing)}")
        try:
            from alibabacloud_docmind_api20220711.client import Client as DocMindClient
            from alibabacloud_tea_openapi import models as openapi_models
        except ImportError as exc:
            raise RuntimeError("缺少阿里云文档解析 SDK 依赖") from exc

        return DocMindClient(
            openapi_models.Config(
                access_key_id=config["access_key_id"],
                access_key_secret=config["access_key_secret"],
                endpoint=DEFAULT_ENDPOINT,
                region_id=str(config.get("region") or DEFAULT_REGION),
                type="access_key",
                connect_timeout=int(self.settings.external_service_timeout_seconds * 1000),
                read_timeout=int(self.settings.external_service_timeout_seconds * 1000),
            )
        )

    def _runtime_options(self):
        from alibabacloud_tea_util import models as util_models

        return util_models.RuntimeOptions(
            connect_timeout=int(self.settings.external_service_timeout_seconds * 1000),
            read_timeout=int(self.settings.external_service_timeout_seconds * 1000),
        )

    def _submit_job(self, path: Path, filename: str, config: dict) -> str:
        try:
            from alibabacloud_docmind_api20220711 import models as docmind_models
        except ImportError as exc:
            raise RuntimeError("缺少阿里云文档解析 SDK 依赖") from exc

        extension = (Path(filename).suffix or path.suffix).lower().lstrip(".")
        if not extension:
            raise bad_request("文档解析需要文件后缀")
        llm_enhancement = _as_bool(config.get("llm_enhancement"), True) or _as_bool(config.get("output_html_table"), False)
        request = docmind_models.SubmitDocParserJobAdvanceRequest(
            file_name=filename,
            file_name_extension=extension,
            formula_enhancement=_as_bool(config.get("formula_enhancement"), False),
            llm_enhancement=llm_enhancement,
            output_html_table=_as_bool(config.get("output_html_table"), False),
            output_format=_output_formats(config.get("output_format")),
            page_index=str(config.get("page_index")).strip() if config.get("page_index") else None,
            enable_event_callback=_as_bool(config.get("enable_event_callback"), False),
        )
        if llm_enhancement and config.get("enhancement_mode", "VLM"):
            request.enhancement_mode = str(config.get("enhancement_mode") or "VLM")
        with path.open("rb") as file_object:
            request.file_url_object = file_object
            response = self._client(config).submit_doc_parser_job_advance(request, self._runtime_options())
        if error_message := _api_error_message(response.body):
            raise bad_request(f"文档解析任务提交失败: {error_message}")
        data = _read_value(response.body, "data") if response.body else None
        task_id = _read_value(data, "id")
        if not task_id:
            body = response.body.to_map() if response.body and hasattr(response.body, "to_map") else {}
            raise bad_request(f"文档解析任务提交失败: {body.get('Message') or body.get('message') or '未返回任务 ID'}")
        return str(task_id)

    def _query_status(self, task_id: str, config: dict) -> dict:
        from alibabacloud_docmind_api20220711 import models as docmind_models

        request = docmind_models.QueryDocParserStatusRequest(id=task_id)
        response = self._client(config).query_doc_parser_status(request)
        if error_message := _api_error_message(response.body):
            raise bad_request(f"文档解析状态查询失败: {error_message}")
        body = _to_plain_data(response.body)
        data = _to_plain_data(_read_value(response.body, "data")) if response.body else None
        if isinstance(data, dict):
            return data
        if isinstance(body, dict):
            return body.get("Data") or body.get("data") or body
        return {}

    @staticmethod
    def _is_retryable_query_error(exc: Exception) -> bool:
        # #61: 仅把"瞬时网络/连接抖动"判为可重试；定性的 API 错误（task 不存在/失败，由
        # _query_status 抛 bad_request）不在此列，避免把瞬时错误等同于需要重新提交付费任务。
        if isinstance(exc, (httpx.TransportError, socket.error, ConnectionError, TimeoutError)):
            return True
        # 阿里云 Tea SDK 的网络层异常名通常含 Connection/Timeout/Throttling 等关键字。
        text = f"{type(exc).__name__} {exc}".lower()
        return any(token in text for token in ("connection", "timed out", "timeout", "temporarily", "throttl", "reset by peer", "network"))

    def _query_status_with_retry(self, task_id: str, config: dict, *, attempts: int = 3, backoff_seconds: float = 1.5) -> dict:
        # #61: 对可重试网络错误做有限重试；耗尽后向上抛出，由调用方决定（绝不静默重新提交）。
        last_exc: Exception | None = None
        for attempt in range(1, attempts + 1):
            try:
                return self._query_status(task_id, config)
            except Exception as exc:  # noqa: BLE001 - 需按可重试性区分处理
                if not self._is_retryable_query_error(exc):
                    raise
                last_exc = exc
                LOGGER.warning(
                    "DocMind status query transient failure (task=%s attempt=%s/%s): %s",
                    task_id,
                    attempt,
                    attempts,
                    exc,
                )
                if attempt < attempts:
                    time.sleep(backoff_seconds * attempt)
        assert last_exc is not None
        raise last_exc

    def _output_format_results(self, status: dict) -> list[dict]:
        items = status.get("OutputFormatResult") or status.get("outputFormatResult") or status.get("output_format_result") or []
        if not isinstance(items, list):
            return []
        return [item for item in items if isinstance(item, dict)]

    def _markdown_result_url(self, status: dict) -> str | None:
        for item in self._output_format_results(status):
            output_type = str(item.get("OutputType") or item.get("outputType") or item.get("output_type") or "").lower()
            url = item.get("OutputFileUrl") or item.get("outputFileUrl") or item.get("output_file_url")
            if output_type == "markdown" and url:
                return str(url)
        return None

    def _wait_for_success(
        self,
        task_id: str,
        config: dict,
        on_progress: Callable[[dict[str, Any]], None] | None = None,
    ) -> dict:
        timeout_seconds = _bounded_int(
            config.get("timeout_seconds") or config.get("timeout"),
            default=DEFAULT_DOC_PARSER_TIMEOUT_SECONDS,
            minimum=30,
            maximum=MAX_DOC_PARSER_TIMEOUT_SECONDS,
        )
        poll_interval = _bounded_int(
            config.get("poll_interval_seconds"),
            default=DEFAULT_DOC_PARSER_POLL_INTERVAL_SECONDS,
            minimum=1,
            maximum=60,
        )
        deadline = time.monotonic() + timeout_seconds
        last_status: dict = {}
        attempt = 0
        while time.monotonic() < deadline:
            attempt += 1
            last_status = self._query_status(task_id, config)
            status = str(last_status.get("Status") or last_status.get("status") or "").lower()
            progress = _safe_int(last_status.get("Processing") or last_status.get("processing"), 0)
            message = last_status.get("Message") or last_status.get("message") or last_status.get("Code")
            if on_progress:
                on_progress(
                    {
                        "stage": "waiting",
                        "docmind_task_id": task_id,
                        "status": status or "unknown",
                        "progress": progress,
                        "attempt": attempt,
                        "timeout_seconds": timeout_seconds,
                        "poll_interval_seconds": poll_interval,
                        "message": str(message) if message else None,
                    }
                )
            if status == "success":
                if on_progress:
                    on_progress(
                        {
                            "stage": "success",
                            "docmind_task_id": task_id,
                            "status": status,
                            "progress": progress or 100,
                            "attempt": attempt,
                            "timeout_seconds": timeout_seconds,
                        }
                    )
                return last_status
            if status in {"fail", "failed"}:
                message = last_status.get("Message") or last_status.get("message") or last_status.get("Code") or "处理失败"
                raise bad_request(f"文档解析失败: {message}")
            time.sleep(poll_interval)
        progress = last_status.get("Processing") or last_status.get("processing") or 0
        if on_progress:
            on_progress(
                {
                    "stage": "timeout",
                    "docmind_task_id": task_id,
                    "status": str(last_status.get("Status") or last_status.get("status") or "timeout").lower(),
                    "progress": _safe_int(progress, 0),
                    "attempt": attempt,
                    "timeout_seconds": timeout_seconds,
                }
            )
        raise bad_request(f"文档解析超时，当前进度 {progress}%")

    def _get_result_batch(self, task_id: str, config: dict, layout_num: int, step_size: int) -> list[dict]:
        from alibabacloud_docmind_api20220711 import models as docmind_models

        request = docmind_models.GetDocParserResultRequest(
            id=task_id,
            layout_num=layout_num,
            layout_step_size=step_size,
        )
        response = self._client(config).get_doc_parser_result(request)
        if error_message := _api_error_message(response.body):
            raise bad_request(f"文档解析结果获取失败: {error_message}")
        data = _to_plain_data(_read_value(response.body, "data")) if response.body else {}
        if isinstance(data, dict):
            layouts = data.get("layouts") or data.get("Layouts") or []
        else:
            layouts = []
        if not isinstance(layouts, list):
            raise bad_request("文档解析结果格式异常")
        return [layout for layout in layouts if isinstance(layout, dict)]

    def _collect_layouts(self, task_id: str, config: dict) -> list[dict]:
        step_size = _bounded_int(
            config.get("layout_step_size"),
            default=DEFAULT_LAYOUT_STEP_SIZE,
            minimum=1,
            maximum=3000,
        )
        layouts: list[dict] = []
        layout_num = 0
        while True:
            batch = self._get_result_batch(task_id, config, layout_num, step_size)
            if not batch:
                break
            layouts.extend(batch)
            layout_num += len(batch)
            if len(batch) < step_size:
                break
        return layouts

    def _download_markdown_result(self, status: dict) -> str:
        url = self._markdown_result_url(status)
        if not url or not _is_remote_url(url):
            return ""
        current_url = url
        response: httpx.Response | None = None
        try:
            with _build_pinned_client(httpx.Timeout(60.0, connect=10.0)) as client:
                for _ in range(MAX_MARKDOWN_IMAGE_REDIRECTS + 1):
                    if not _is_remote_url(current_url):
                        return ""
                    response = client.get(current_url)
                    if bool(getattr(response, "is_redirect", False)):
                        location = response.headers.get("location")
                        if not location:
                            return ""
                        # 用原始 current_url（非被改写成 IP 的 response.url）解析相对重定向，
                        # 下一跳仍按真实域名重新解析+校验并锁 IP。
                        current_url = str(httpx.URL(current_url).join(location))
                        continue
                    response.raise_for_status()
                    break
                else:
                    return ""
        except httpx.HTTPError:
            return ""
        if response is None:
            return ""
        content = b""
        if hasattr(response, "iter_bytes"):
            for chunk in response.iter_bytes(chunk_size=64 * 1024):
                content += chunk
                if len(content) > MAX_MARKDOWN_RESULT_BYTES:
                    return ""
        else:
            content = getattr(response, "content", b"")
        if not content or len(content) > MAX_MARKDOWN_RESULT_BYTES:
            return ""
        return content.decode(response.encoding or "utf-8", errors="ignore").strip()

    def _pages_from_successful_task(self, task_id: str, config: dict, status: dict, db: Session | None) -> list[dict]:
        pages = self._pages_from_layouts(self._collect_layouts(task_id, config))
        markdown = self._download_markdown_result(status)
        if markdown:
            markdown_pages = _pages_from_markdown(markdown)
            layout_text_len = sum(len(page.get("page_text") or "") for page in pages)
            markdown_text_len = sum(len(page.get("page_text") or "") for page in markdown_pages)
            if markdown_pages and (not pages or (len(markdown_pages) >= len(pages) and markdown_text_len > layout_text_len)):
                pages = markdown_pages
        image_cache: dict[str, str | None] = {}
        for page in pages:
            page["page_text"] = _localize_markdown_images(page.get("page_text") or "", db, image_cache)
        return pages

    def _pages_from_layouts(self, layouts: list[dict]) -> list[dict]:
        grouped: dict[int, list[dict]] = defaultdict(list)
        for index, layout in enumerate(layouts):
            raw_page = layout.get("pageNum", layout.get("page_num", layout.get("pageNumber", 0)))
            try:
                page_number = int(raw_page) + 1
            except (TypeError, ValueError):
                page_number = 1
            layout["_fallback_order"] = index
            grouped[page_number].append(layout)

        pages: list[dict] = []
        for page_number in sorted(grouped):
            page_layouts = sorted(grouped[page_number], key=lambda item: _safe_int(item.get("index"), item["_fallback_order"]))
            pieces = [piece for piece in (_layout_text(layout) for layout in page_layouts) if piece]
            title = None
            for layout in page_layouts:
                if str(layout.get("type", "")).lower() == "title":
                    title = _layout_title(layout)
                    break
            pages.append(_normalize_page(title or f"第{page_number}页", "\n\n".join(pieces), page_number))
        return pages

    def parse(
        self,
        path: Path,
        material_type: str,
        db: Session | None,
        filename: str | None = None,
        resume_task_id: str | None = None,
        on_progress: Callable[[dict[str, Any]], None] | None = None,
    ) -> list[dict]:
        if material_type not in SUPPORTED_MATERIAL_TYPES:
            raise bad_request("暂不支持该资料类型解析")
        # 修复 DEF-02：图片资料不走 DocMind，改由 OCR 抽取文本生成课时页。
        if material_type == MaterialType.IMAGE.value:
            return _parse_image(path, db)
        service = get_enabled_service_config(db, SERVICE_TYPE)
        if service is None:
            if self.settings.app_env != "production":
                return _parse_local_fallback(path, material_type)
            raise bad_request("文档解析服务未配置，请先在管理员后台配置阿里云文档解析服务")
        if service.provider != "aliyun":
            raise bad_request(f"暂不支持的文档解析服务提供方: {service.provider}")

        if resume_task_id:
            try:
                # #61: 带有限重试的状态查询，区分"任务确实不存在/已失败"与"瞬时网络抖动"。
                status = self._query_status_with_retry(resume_task_id, service.config)
                if str(status.get("Status") or status.get("status") or "").lower() == "success":
                    if on_progress:
                        on_progress({"stage": "resuming", "docmind_task_id": resume_task_id, "status": "success", "progress": 100})
                    pages = self._pages_from_successful_task(resume_task_id, service.config, status, db)
                    if pages:
                        if on_progress:
                            on_progress(
                                {
                                    "stage": "parsed",
                                    "docmind_task_id": resume_task_id,
                                    "status": "success",
                                    "progress": 100,
                                    "page_count": len(pages),
                                    "resumed": True,
                                }
                            )
                        return pages
            except Exception as exc:
                if on_progress:
                    on_progress(
                        {
                            "stage": "resume_failed",
                            "docmind_task_id": resume_task_id,
                            "status": "failed",
                            "progress": 0,
                            "message": str(exc),
                        }
                    )
                # #61: 若耗尽重试后仍是可重试的瞬时网络错误，不能把它当成"任务不存在/失败"而
                # 静默重新提交（会重复触发付费 DocMind 任务）。直接上抛友好错误，让上层稍后重试。
                if self._is_retryable_query_error(exc):
                    LOGGER.warning(
                        "DocMind resume status query failed transiently after retries (task=%s): %s",
                        resume_task_id,
                        exc,
                    )
                    raise bad_request("文档解析状态查询暂时不可用，请稍后重试") from exc
                # 其余（任务确实不存在/已失败等定性错误）：按原逻辑回退到重新提交。
                LOGGER.info(
                    "DocMind resume task unavailable, will resubmit (task=%s): %s",
                    resume_task_id,
                    exc,
                )

        task_id = self._submit_job(path, filename or path.name, service.config)
        if on_progress:
            on_progress({"stage": "submitted", "docmind_task_id": task_id, "status": "submitted", "progress": 0})
        status = self._wait_for_success(task_id, service.config, on_progress=on_progress)
        if on_progress:
            on_progress({"stage": "collecting", "docmind_task_id": task_id, "status": "success", "progress": 100})
        pages = self._pages_from_successful_task(task_id, service.config, status, db)
        if not pages:
            raise bad_request("文档解析未返回有效内容")
        if on_progress:
            on_progress(
                {
                    "stage": "parsed",
                    "docmind_task_id": task_id,
                    "status": "success",
                    "progress": 100,
                    "page_count": len(pages),
                }
            )
        return pages


doc_parser_service = DocParserService()


def parse_material(
    path: Path,
    material_type: str,
    db: Session | None = None,
    filename: str | None = None,
    resume_task_id: str | None = None,
    on_progress: Callable[[dict[str, Any]], None] | None = None,
) -> list[dict]:
    return doc_parser_service.parse(path, material_type, db, filename=filename, resume_task_id=resume_task_id, on_progress=on_progress)
