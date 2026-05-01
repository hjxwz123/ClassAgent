from io import BytesIO

from pptx import Presentation


def register_user(client, *, email, password, nickname, role, student_no=None, employee_no=None):
    payload = {
        "email": email,
        "password": password,
        "nickname": nickname,
        "role": role,
        "student_no": student_no,
        "employee_no": employee_no,
    }
    if role == "teacher":
        admin_login = login_user(client, email="admin@classagent.com", password="Admin123456")
        response = client.post("/api/v1/admin/users/admin", json=payload, headers=auth_headers(admin_login["access_token"]))
    else:
        response = client.post("/api/v1/auth/register", json=payload)
    assert response.status_code == 200, response.text


def login_user(client, *, email, password):
    response = client.post("/api/v1/auth/login", json={"email": email, "password": password})
    assert response.status_code == 200, response.text
    return response.json()["data"]


def auth_headers(token: str) -> dict[str, str]:
    return {"Authorization": f"Bearer {token}"}


def build_pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "特征值"
    slide.placeholders[1].text = "特征值描述线性变换的固有缩放因子。"
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def prepare_activity_data(client):
    register_user(
        client,
        email="teacher4@example.com",
        password="Teacher123",
        nickname="吴老师",
        role="teacher",
        employee_no="T2026004",
    )
    register_user(
        client,
        email="teacher5@example.com",
        password="Teacher123",
        nickname="郑老师",
        role="teacher",
        employee_no="T2026005",
    )
    register_user(
        client,
        email="student4@example.com",
        password="Student123",
        nickname="孙同学",
        role="student",
        student_no="S2026004",
    )
    teacher_login = login_user(client, email="teacher4@example.com", password="Teacher123")
    teacher_headers = auth_headers(teacher_login["access_token"])
    student_login = login_user(client, email="student4@example.com", password="Student123")
    student_headers = auth_headers(student_login["access_token"])

    course_resp = client.post(
        "/api/v1/courses",
        json={"name": "高阶线代", "description": "特征值与特征向量", "term": "2026春"},
        headers=teacher_headers,
    )
    course = course_resp.json()["data"]
    chapter_resp = client.post(
        f"/api/v1/courses/{course['id']}/chapters",
        json={"title": "特征值", "description": "核心章节", "order_index": 1},
        headers=teacher_headers,
    )
    chapter = chapter_resp.json()["data"]
    client.post("/api/v1/courses/join", json={"course_code": course["course_code"]}, headers=student_headers)
    upload_resp = client.post(
        "/api/v1/materials",
        data={
            "course_id": str(course["id"]),
            "title": "特征值课件",
            "category": "courseware",
            "chapter_id": str(chapter["id"]),
        },
        files={
            "file": (
                "eigen.pptx",
                build_pptx_bytes(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
        headers=teacher_headers,
    )
    material = upload_resp.json()["data"]
    detail_resp = client.get(f"/api/v1/materials/{material['id']}", headers=teacher_headers)
    lesson_id = detail_resp.json()["data"]["lesson_id"]
    client.post(f"/api/v1/lessons/{lesson_id}/publish", headers=teacher_headers)
    client.post(
        f"/api/v1/lessons/{lesson_id}/progress",
        json={"current_page": 1, "added_seconds": 180, "completed": True},
        headers=student_headers,
    )
    client.post(
        "/api/v1/qa/ask",
        json={"course_id": course["id"], "question": "什么是特征值"},
        headers=student_headers,
    )
    client.post(
        "/api/v1/tutoring/problems/text",
        json={"course_id": course["id"], "text": "已知矩阵A，求特征值"},
        headers=student_headers,
    )
    quiz_resp = client.post(
        "/api/v1/learning/quizzes/generate",
        json={"course_id": course["id"], "chapter_id": chapter["id"], "title": "特征值测验", "quiz_type": "course", "question_count": 2},
        headers=teacher_headers,
    )
    quiz = quiz_resp.json()["data"]
    client.post(f"/api/v1/learning/quizzes/{quiz['id']}/publish", headers=teacher_headers)
    quiz_detail = client.get(f"/api/v1/learning/quizzes/{quiz['id']}", headers=student_headers).json()["data"]
    questions = quiz_detail["questions"]
    client.post(
        f"/api/v1/learning/quizzes/{quiz['id']}/submit",
        json={"answers": [{"question_id": questions[0]["id"], "answer": 0}, {"question_id": questions[1]["id"], "answer": 0}]},
        headers=student_headers,
    )
    return course, teacher_headers


def test_teacher_analytics_and_admin_operations(client):
    course, teacher_headers = prepare_activity_data(client)
    admin_login = login_user(client, email="admin@classagent.com", password="Admin123456")
    admin_headers = auth_headers(admin_login["access_token"])

    dashboard_resp = client.get("/api/v1/admin/dashboard", params={"activity_days": 7}, headers=admin_headers)
    assert dashboard_resp.status_code == 200, dashboard_resp.text
    assert len(dashboard_resp.json()["data"]["activity_trend"]) == 7

    analytics_resp = client.get(f"/api/v1/analytics/courses/{course['id']}", headers=teacher_headers)
    assert analytics_resp.status_code == 200, analytics_resp.text
    analytics = analytics_resp.json()["data"]
    assert "high_frequency_questions" in analytics
    assert "suggestion" in analytics

    users_resp = client.get("/api/v1/admin/users", headers=admin_headers)
    assert users_resp.status_code == 200, users_resp.text
    assert len(users_resp.json()["data"]) >= 4
    target_student = next(item for item in users_resp.json()["data"] if item["email"] == "student4@example.com")
    target_teacher2 = next(item for item in users_resp.json()["data"] if item["email"] == "teacher5@example.com")

    create_admin_resp = client.post(
        "/api/v1/admin/users/admin",
        json={"email": "admin2@example.com", "password": "Admin123456", "nickname": "二号管理员"},
        headers=admin_headers,
    )
    assert create_admin_resp.status_code == 200, create_admin_resp.text

    update_user_resp = client.patch(
        f"/api/v1/admin/users/{target_student['id']}",
        json={"status": "disabled"},
        headers=admin_headers,
    )
    assert update_user_resp.status_code == 200, update_user_resp.text
    assert update_user_resp.json()["data"]["status"] == "disabled"

    reset_pwd_resp = client.post(
        f"/api/v1/admin/users/{target_student['id']}/reset-password",
        json={"new_password": "Student999"},
        headers=admin_headers,
    )
    assert reset_pwd_resp.status_code == 200, reset_pwd_resp.text

    courses_resp = client.get("/api/v1/admin/courses", headers=admin_headers)
    assert courses_resp.status_code == 200, courses_resp.text
    assert len(courses_resp.json()["data"]) >= 1

    course_detail_resp = client.get(f"/api/v1/admin/courses/{course['id']}", headers=admin_headers)
    assert course_detail_resp.status_code == 200, course_detail_resp.text
    assert course_detail_resp.json()["data"]["material_count"] >= 1

    takeover_resp = client.post(
        f"/api/v1/admin/courses/{course['id']}/takeover",
        json={"teacher_id": target_teacher2["id"]},
        headers=admin_headers,
    )
    assert takeover_resp.status_code == 200, takeover_resp.text
    assert takeover_resp.json()["data"]["teacher_id"] == target_teacher2["id"]

    materials_resp = client.get("/api/v1/admin/materials", headers=admin_headers)
    assert materials_resp.status_code == 200, materials_resp.text
    assert len(materials_resp.json()["data"]) >= 1
    material_id = materials_resp.json()["data"][0]["id"]
    filtered_materials_resp = client.get(
        "/api/v1/admin/materials",
        params={"material_type": "pptx", "teacher_id": materials_resp.json()["data"][0]["uploader_id"]},
        headers=admin_headers,
    )
    assert filtered_materials_resp.status_code == 200, filtered_materials_resp.text
    assert len(filtered_materials_resp.json()["data"]) >= 1

    material_stats_resp = client.get("/api/v1/admin/materials/stats", headers=admin_headers)
    assert material_stats_resp.status_code == 200, material_stats_resp.text
    assert material_stats_resp.json()["data"]["total"] >= 1
    assert material_stats_resp.json()["data"]["by_type"]

    model_save_resp = client.post(
        "/api/v1/admin/model-configs",
        json={
            "provider": "mock",
            "model_name": "mock-v1",
            "purpose": "qa",
            "endpoint": None,
            "api_key": "mock-key",
            "is_default": True,
            "extra_config": {"note": "test"},
        },
        headers=admin_headers,
    )
    assert model_save_resp.status_code == 200, model_save_resp.text
    model_id = model_save_resp.json()["data"]["id"]

    model_test_resp = client.post(f"/api/v1/admin/model-configs/{model_id}/test", headers=admin_headers)
    assert model_test_resp.status_code == 200, model_test_resp.text
    assert model_test_resp.json()["data"]["success"] is True

    embedding_save_resp = client.post(
        "/api/v1/admin/model-configs",
        json={
            "provider": "mock",
            "model_name": "mock-embedding",
            "purpose": "embedding",
            "endpoint": None,
            "api_key": "mock-key",
            "is_default": True,
            "extra_config": {"dimensions": 384},
        },
        headers=admin_headers,
    )
    assert embedding_save_resp.status_code == 200, embedding_save_resp.text
    embedding_id = embedding_save_resp.json()["data"]["id"]
    embedding_test_resp = client.post(f"/api/v1/admin/model-configs/{embedding_id}/test", headers=admin_headers)
    assert embedding_test_resp.status_code == 200, embedding_test_resp.text
    assert embedding_test_resp.json()["data"]["success"] is True

    service_save_resp = client.post(
        "/api/v1/admin/service-configs",
        json={
            "service_type": "tts",
            "provider": "mock",
            "name": "mock-tts",
            "config": {"voice": "xiaoyun"},
            "is_enabled": True,
        },
        headers=admin_headers,
    )
    assert service_save_resp.status_code == 200, service_save_resp.text
    service_id = service_save_resp.json()["data"]["id"]

    service_test_resp = client.post(f"/api/v1/admin/service-configs/{service_id}/test", headers=admin_headers)
    assert service_test_resp.status_code == 200, service_test_resp.text
    assert service_test_resp.json()["data"]["success"] is True

    doc_parser_save_resp = client.post(
        "/api/v1/admin/service-configs",
        json={
            "service_type": "doc_parser",
            "provider": "aliyun",
            "name": "aliyun-doc-parser",
            "config": {
                "access_key_id": "test-ak",
                "access_key_secret": "test-secret",
                "region": "cn-hangzhou",
            },
            "is_enabled": True,
        },
        headers=admin_headers,
    )
    assert doc_parser_save_resp.status_code == 200, doc_parser_save_resp.text
    doc_parser_id = doc_parser_save_resp.json()["data"]["id"]
    doc_parser_test_resp = client.post(f"/api/v1/admin/service-configs/{doc_parser_id}/test", headers=admin_headers)
    assert doc_parser_test_resp.status_code == 200, doc_parser_test_resp.text
    assert doc_parser_test_resp.json()["data"]["success"] is True

    email_save_resp = client.post(
        "/api/v1/admin/service-configs",
        json={
            "service_type": "email",
            "provider": "mock",
            "name": "mock-email",
            "config": {"host": "localhost", "port": 25, "sender": "noreply@example.com"},
            "is_enabled": True,
        },
        headers=admin_headers,
    )
    assert email_save_resp.status_code == 200, email_save_resp.text
    email_id = email_save_resp.json()["data"]["id"]
    email_test_resp = client.post(f"/api/v1/admin/service-configs/{email_id}/test", headers=admin_headers)
    assert email_test_resp.status_code == 200, email_test_resp.text
    assert email_test_resp.json()["data"]["success"] is True

    settings_resp = client.get("/api/v1/admin/system-settings", headers=admin_headers)
    assert settings_resp.status_code == 200, settings_resp.text
    assert len(settings_resp.json()["data"]) >= 1

    update_setting_resp = client.put(
        "/api/v1/admin/system-settings/system.announcement",
        json={"value": "期中周系统维护公告"},
        headers=admin_headers,
    )
    assert update_setting_resp.status_code == 200, update_setting_resp.text

    monitoring_resp = client.get("/api/v1/admin/monitoring/overview", headers=admin_headers)
    assert monitoring_resp.status_code == 200, monitoring_resp.text
    monitoring = monitoring_resp.json()["data"]
    assert monitoring["api_call_count_30m"] > 0

    model_usage_resp = client.get("/api/v1/admin/model-usage", headers=admin_headers)
    assert model_usage_resp.status_code == 200, model_usage_resp.text
    assert len(model_usage_resp.json()["data"]["items"]) >= 1

    login_logs_resp = client.get("/api/v1/admin/logs/login", headers=admin_headers)
    assert login_logs_resp.status_code == 200, login_logs_resp.text
    login_logs_filtered_resp = client.get(
        "/api/v1/admin/logs/login",
        params={"user_id": target_student["id"], "limit": 10},
        headers=admin_headers,
    )
    assert login_logs_filtered_resp.status_code == 200, login_logs_filtered_resp.text
    operation_logs_resp = client.get("/api/v1/admin/logs/operations", headers=admin_headers)
    assert operation_logs_resp.status_code == 200, operation_logs_resp.text
    operation_logs_filtered_resp = client.get(
        "/api/v1/admin/logs/operations",
        params={"action": "material.create", "target_type": "material", "limit": 10},
        headers=admin_headers,
    )
    assert operation_logs_filtered_resp.status_code == 200, operation_logs_filtered_resp.text
    error_logs_resp = client.get("/api/v1/admin/logs/errors", headers=admin_headers)
    assert error_logs_resp.status_code == 200, error_logs_resp.text
    error_logs_filtered_resp = client.get(
        "/api/v1/admin/logs/errors",
        params={"level": "error", "limit": 10},
        headers=admin_headers,
    )
    assert error_logs_filtered_resp.status_code == 200, error_logs_filtered_resp.text

    backup_create_resp = client.post("/api/v1/admin/backups", headers=admin_headers)
    assert backup_create_resp.status_code == 200, backup_create_resp.text
    assert backup_create_resp.json()["data"]["status"] == "success"
    backup_id = backup_create_resp.json()["data"]["id"]
    backup_list_resp = client.get("/api/v1/admin/backups", headers=admin_headers)
    assert backup_list_resp.status_code == 200, backup_list_resp.text
    assert len(backup_list_resp.json()["data"]) >= 1

    delete_material_resp = client.delete(f"/api/v1/admin/materials/{material_id}", headers=admin_headers)
    assert delete_material_resp.status_code == 200, delete_material_resp.text

    backup_restore_resp = client.post(f"/api/v1/admin/backups/{backup_id}/restore", headers=admin_headers)
    assert backup_restore_resp.status_code == 200, backup_restore_resp.text
    backup_delete_resp = client.delete(f"/api/v1/admin/backups/{backup_id}", headers=admin_headers)
    assert backup_delete_resp.status_code == 200, backup_delete_resp.text
