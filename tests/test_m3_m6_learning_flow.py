import json
from io import BytesIO

import pytest
from sqlalchemy import delete, select
from pptx import Presentation

from app.core.enums import LessonStatus, MaterialCategory, MaterialType, ProcessStatus
from app.core.errors import AppError
from app.db import session as db_session
from app.db.models import (
    Chapter,
    Course,
    CourseMaterial,
    KnowledgeChunk,
    KnowledgePoint,
    Lesson,
    LessonPage,
    ProblemGuidance,
    QuizQuestion,
    StudentLearningSignal,
    User,
)
from app.services.ai import _pack_rag_contexts, ai_service, sanitize_quiz_source_text
from app.services.learning import QUIZ_SOURCE_CONTEXT_HARD_LIMIT, _course_source_text_for_quiz
from app.services.pedagogy import ensure_lesson_pedagogy_artifacts, page_activity_payload
from tests.auth_helpers import request_registration_token


PNG_BYTES = (
    b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
    b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\rIDATx\x9cc\xf8\xcf\xc0"
    b"\xf0\x1f\x00\x05\x05\x02\x00\x1e^\x99\xed\x00\x00\x00\x00IEND\xaeB`\x82"
)


def fake_quiz_questions(*, topic, source_text, count, db=None, **kwargs):
    items = [
        {
            "question_type": "single_choice",
            "stem": "矩阵可以表示线性变换，下列说法哪一项正确？",
            "options": ["矩阵可以表示线性变换", "矩阵只能表示图片", "矩阵与线性变换无关", "矩阵不需要理解条件"],
            "reference_answer": {"value": 0},
            "explanation": "课程资料明确提到矩阵可以表示线性变换。",
            "score": 10,
            "difficulty": "standard",
        },
        {
            "question_type": "judge",
            "stem": "判断：行列式可以反映线性变换中的缩放系数。",
            "options": ["正确", "错误"],
            "reference_answer": {"value": 0},
            "explanation": "课程资料提到行列式反映缩放系数。",
            "score": 10,
            "difficulty": "standard",
        },
        {
            "question_type": "short_answer",
            "stem": "请简述矩阵与线性变换之间的关系。",
            "options": None,
            "reference_answer": {"keywords": ["矩阵", "线性", "变换"]},
            "explanation": "应围绕矩阵可以表示线性变换展开。",
            "score": 20,
            "difficulty": "advanced",
        },
    ]
    return items[:count]


def fake_quiz_questions_with_types(*, topic, source_text, count, type_counts=None, db=None, **kwargs):
    typed_items = {
        "single_choice": {
            "question_type": "single_choice",
            "stem": f"{topic} 中，矩阵可以表示什么？",
            "options": ["线性变换", "图片文件名", "课程编号", "无关文本"],
            "reference_answer": {"value": 0},
            "explanation": "课程资料说明矩阵可以表示线性变换。",
            "score": 10,
            "difficulty": "standard",
        },
        "judge": {
            "question_type": "judge",
            "stem": f"判断：{topic} 与线性变换有关。",
            "options": ["正确", "错误"],
            "reference_answer": {"value": 0},
            "explanation": "课程资料围绕矩阵和线性变换展开。",
            "score": 10,
            "difficulty": "standard",
        },
        "blank": {
            "question_type": "blank",
            "stem": f"{topic} 可以表示____。",
            "options": None,
            "reference_answer": {"keywords": ["线性变换"]},
            "explanation": "填入线性变换。",
            "score": 10,
            "difficulty": "standard",
        },
        "short_answer": {
            "question_type": "short_answer",
            "stem": f"请简述 {topic} 与线性变换的关系。",
            "options": None,
            "reference_answer": {"keywords": ["矩阵", "线性", "变换"]},
            "explanation": "应围绕矩阵表示线性变换展开。",
            "score": 20,
            "difficulty": "advanced",
        },
    }
    if type_counts:
        items = []
        for question_type, type_count in type_counts.items():
            for index in range(int(type_count)):
                item = {**typed_items[question_type]}
                item["stem"] = f"{item['stem']}（{index + 1}）"
                items.append(item)
        return items[:count]
    return fake_quiz_questions(topic=topic, source_text=source_text, count=count, db=db)


def register_user(client, *, email, password, nickname, role, student_no=None, employee_no=None):
    payload = {
        "email": email,
        "password": password,
        "nickname": nickname,
        "role": role,
        "student_no": student_no,
        "employee_no": employee_no,
    }
    if role == "teacher":
        admin_login = login_user(client, email="admin@classagent.com", password="Admin123456")
        response = client.post("/api/v1/admin/users/admin", json=payload, headers=auth_headers(admin_login["access_token"]))
    else:
        payload["token"] = request_registration_token(client, email)
        response = client.post("/api/v1/auth/register", json=payload)
    assert response.status_code == 200, response.text


def login_user(client, *, email, password):
    response = client.post("/api/v1/auth/login", json={"email": email, "password": password})
    assert response.status_code == 200, response.text
    return response.json()["data"]


def auth_headers(token: str) -> dict[str, str]:
    return {"Authorization": f"Bearer {token}"}


def build_pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "矩阵"
    slide.placeholders[1].text = "矩阵可以表示线性变换。"
    slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide2.shapes.title.text = "行列式"
    slide2.placeholders[1].text = "行列式反映缩放系数。"
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def bootstrap_course_with_material(client):
    register_user(
        client,
        email="teacher3@example.com",
        password="Teacher123",
        nickname="周老师",
        role="teacher",
        employee_no="T2026003",
    )
    register_user(
        client,
        email="student3@example.com",
        password="Student123",
        nickname="钱同学",
        role="student",
        student_no="S2026003",
    )
    teacher_login = login_user(client, email="teacher3@example.com", password="Teacher123")
    student_login = login_user(client, email="student3@example.com", password="Student123")
    teacher_headers = auth_headers(teacher_login["access_token"])
    student_headers = auth_headers(student_login["access_token"])

    course_resp = client.post(
        "/api/v1/courses",
        json={"name": "线性代数进阶", "description": "矩阵、向量与行列式", "term": "2026春"},
        headers=teacher_headers,
    )
    course = course_resp.json()["data"]
    chapter_resp = client.post(
        f"/api/v1/courses/{course['id']}/chapters",
        json={"title": "第一章 矩阵基础", "description": "矩阵概念", "order_index": 1},
        headers=teacher_headers,
    )
    chapter = chapter_resp.json()["data"]
    join_resp = client.post(
        "/api/v1/courses/join",
        json={"course_code": course["course_code"]},
        headers=student_headers,
    )
    assert join_resp.status_code == 200, join_resp.text

    upload_resp = client.post(
        "/api/v1/materials",
        data={
            "course_id": str(course["id"]),
            "title": "矩阵与行列式课件",
            "category": "courseware",
            "chapter_id": str(chapter["id"]),
        },
        files={
            "file": (
                "matrix_lesson.pptx",
                build_pptx_bytes(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
        headers=teacher_headers,
    )
    assert upload_resp.status_code == 200, upload_resp.text
    material = upload_resp.json()["data"]
    detail_resp = client.get(f"/api/v1/materials/{material['id']}", headers=teacher_headers)
    assert detail_resp.status_code == 200, detail_resp.text
    lesson_id = detail_resp.json()["data"]["lesson_id"]
    publish_resp = client.post(f"/api/v1/lessons/{lesson_id}/publish", headers=teacher_headers)
    assert publish_resp.status_code == 200, publish_resp.text
    return course, chapter, lesson_id, teacher_headers, student_headers


def test_quiz_generation_calls_ai_with_course_context_without_material(client, monkeypatch):
    register_user(
        client,
        email="teacher-no-material@example.com",
        password="Teacher123",
        nickname="无资料老师",
        role="teacher",
        employee_no="T2026099",
    )
    teacher_login = login_user(client, email="teacher-no-material@example.com", password="Teacher123")
    teacher_headers = auth_headers(teacher_login["access_token"])
    course_resp = client.post(
        "/api/v1/courses",
        json={"name": "离散数学", "description": "图论与集合基础", "term": "2026春"},
        headers=teacher_headers,
    )
    assert course_resp.status_code == 200, course_resp.text
    course = course_resp.json()["data"]
    captured = {}

    def fake_ai_quiz_questions(*, topic, source_text, count, db=None, **kwargs):
        captured["topic"] = topic
        captured["source_text"] = source_text
        return fake_quiz_questions(topic=topic, source_text=source_text, count=count, db=db)

    monkeypatch.setattr(ai_service, "generate_quiz_questions", fake_ai_quiz_questions)
    quiz_resp = client.post(
        "/api/v1/learning/quizzes/generate",
        json={"course_id": course["id"], "title": "课程基础测验", "quiz_type": "course", "question_count": 1},
        headers=teacher_headers,
    )
    assert quiz_resp.status_code == 200, quiz_resp.text
    assert captured["topic"] == "离散数学"
    assert "课程名称：离散数学" in captured["source_text"]
    assert "课程简介：图论与集合基础" in captured["source_text"]


def test_quiz_generation_never_extracts_knowledge_points_on_user_request(client, monkeypatch):
    register_user(
        client,
        email="teacher-no-lazy-points@example.com",
        password="Teacher123",
        nickname="不懒加载老师",
        role="teacher",
        employee_no="T2026098",
    )
    teacher_login = login_user(client, email="teacher-no-lazy-points@example.com", password="Teacher123")
    teacher_headers = auth_headers(teacher_login["access_token"])
    course_resp = client.post(
        "/api/v1/courses",
        json={"name": "编译原理热路径", "description": "词法分析与语法分析", "term": "2026春"},
        headers=teacher_headers,
    )
    assert course_resp.status_code == 200, course_resp.text
    course = course_resp.json()["data"]
    chapter_resp = client.post(
        f"/api/v1/courses/{course['id']}/chapters",
        json={"title": "词法分析", "description": "", "order_index": 1},
        headers=teacher_headers,
    )
    assert chapter_resp.status_code == 200, chapter_resp.text
    chapter = chapter_resp.json()["data"]

    with db_session.SessionLocal() as db:
        db.add(
            KnowledgeChunk(
                course_id=course["id"],
                chapter_id=chapter["id"],
                title="词法分析片段",
                content="词法分析器将源程序字符流转换为单词符号序列。",
            )
        )
        db.commit()

    def fail_lazy_extraction(*args, **kwargs):
        raise AssertionError("quiz generation must not extract knowledge points on the user request path")

    monkeypatch.setattr(ai_service, "extract_knowledge_points", fail_lazy_extraction)
    monkeypatch.setattr(ai_service, "generate_quiz_questions", fake_quiz_questions)
    response = client.post(
        "/api/v1/learning/quizzes/generate",
        json={
            "course_id": course["id"],
            "chapter_id": chapter["id"],
            "title": "词法分析测验",
            "quiz_type": "course",
            "question_count": 1,
        },
        headers=teacher_headers,
    )
    assert response.status_code == 200, response.text
    assert response.json()["data"]["id"] > 0


def test_lesson_detail_includes_material_for_original_preview(client):
    _course, _chapter, lesson_id, _teacher_headers, student_headers = bootstrap_course_with_material(client)

    lesson_detail_resp = client.get(f"/api/v1/lessons/{lesson_id}", headers=student_headers)
    assert lesson_detail_resp.status_code == 200, lesson_detail_resp.text
    detail = lesson_detail_resp.json()["data"]
    assert detail["lesson"]["id"] == lesson_id
    assert detail["material"]["id"] == detail["lesson"]["material_id"]
    assert detail["material"]["material_type"] == "pptx"
    assert detail["material"]["preview_url"] is None
    assert len(detail["pages"]) >= 1


def test_quiz_source_context_covers_many_pages_without_oversized_prompt(client):
    teacher_email = "teacher-long-ppt@example.com"
    page_count = 420
    register_user(
        client,
        email=teacher_email,
        password="Teacher123",
        nickname="长课件老师",
        role="teacher",
        employee_no="T2026100",
    )
    teacher_login = login_user(client, email=teacher_email, password="Teacher123")
    teacher_headers = auth_headers(teacher_login["access_token"])
    course_resp = client.post(
        "/api/v1/courses",
        json={"name": "编译原理", "description": "属性文法与语义分析", "term": "2026春"},
        headers=teacher_headers,
    )
    assert course_resp.status_code == 200, course_resp.text
    course = course_resp.json()["data"]
    chapter_resp = client.post(
        f"/api/v1/courses/{course['id']}/chapters",
        json={"title": "属性文法", "description": "", "order_index": 1},
        headers=teacher_headers,
    )
    assert chapter_resp.status_code == 200, chapter_resp.text
    chapter = chapter_resp.json()["data"]

    with db_session.SessionLocal() as db:
        teacher = db.scalar(select(User).where(User.email == teacher_email))
        assert teacher is not None
        material = CourseMaterial(
            course_id=course["id"],
            chapter_id=chapter["id"],
            uploader_id=teacher.id,
            title="属性文法长课件",
            category=MaterialCategory.COURSEWARE.value,
            material_type=MaterialType.PPTX.value,
            size_bytes=8192,
            original_filename="attribute-grammar.pptx",
            storage_path="uploads/test/attribute-grammar.pptx",
            preview_url="/static/uploads/test/attribute-grammar.pptx",
            extracted_text=None,
            parse_status=ProcessStatus.READY.value,
            vector_status=ProcessStatus.READY.value,
        )
        db.add(material)
        db.commit()
        db.refresh(material)
        lesson = Lesson(
            course_id=course["id"],
            chapter_id=chapter["id"],
            material_id=material.id,
            title="属性文法长课件",
            summary="长 PPT 每页覆盖",
            page_count=page_count,
            status=LessonStatus.PUBLISHED.value,
        )
        db.add(lesson)
        db.commit()
        db.refresh(lesson)
        db.add_all(
            [
                LessonPage(
                    lesson_id=lesson.id,
                    page_number=index,
                    page_title=f"第{index}页主题",
                    page_text=(
                        f"第{index}页开头知识点：属性文法页首覆盖标记{index:03d}。"
                        + "本页围绕属性文法、语义规则、综合属性和继承属性展开说明。" * 30
                        + f"第{index}页结尾知识点：属性文法页尾覆盖标记{index:03d}。"
                    ),
                    script_status=ProcessStatus.READY.value,
                )
                for index in range(1, page_count + 1)
            ]
        )
        db.commit()

        source_text = _course_source_text_for_quiz(db, course_id=course["id"], chapter_ids=[chapter["id"]], points=[])

    assert len(source_text) <= QUIZ_SOURCE_CONTEXT_HARD_LIMIT
    # 新语义：超限时按相关性整片取舍——保留的页必须【完整】(页首/页尾标记成对出现，
    # 不再被头/中/尾采样拦腰打碎)，未入选的页整片舍弃。这里断言：入选页均完整、
    # 且在 80k 预算内覆盖了足量页面。
    kept_pages = [
        index
        for index in range(1, page_count + 1)
        if f"属性文法页首覆盖标记{index:03d}" in source_text
    ]
    assert len(kept_pages) >= 40, f"仅保留 {len(kept_pages)} 页，预算内应能容纳更多完整页"
    for index in kept_pages:
        assert f"属性文法页尾覆盖标记{index:03d}" in source_text, f"第{index}页被截断，整片取舍语义被破坏"
    assert " ... " not in source_text, "出现头/中/尾采样分隔符，片段被打碎"


def test_learning_core_flow(client, monkeypatch):
    monkeypatch.setattr(ai_service, "generate_quiz_questions", fake_quiz_questions)
    course, chapter, lesson_id, teacher_headers, student_headers = bootstrap_course_with_material(client)

    lessons_resp = client.get("/api/v1/lessons", params={"course_id": course["id"]}, headers=student_headers)
    assert lessons_resp.status_code == 200, lessons_resp.text
    assert len(lessons_resp.json()["data"]) == 1

    lesson_detail_resp = client.get(f"/api/v1/lessons/{lesson_id}", headers=student_headers)
    assert lesson_detail_resp.status_code == 200, lesson_detail_resp.text
    lesson_detail = lesson_detail_resp.json()["data"]
    assert len(lesson_detail["pages"]) == 2
    first_page_id = lesson_detail["pages"][0]["id"]

    progress_resp = client.post(
        f"/api/v1/lessons/{lesson_id}/progress",
        json={"current_page": 1, "added_seconds": 120, "completed": False},
        headers=student_headers,
    )
    assert progress_resp.status_code == 200, progress_resp.text
    assert progress_resp.json()["data"]["current_page"] == 1

    qa_resp = client.post(
        "/api/v1/qa/ask",
        json={"course_id": course["id"], "lesson_page_id": first_page_id, "question": "矩阵可以表示什么"},
        headers=student_headers,
    )
    assert qa_resp.status_code == 200, qa_resp.text
    qa_data = qa_resp.json()["data"]
    assert qa_data["is_out_of_scope"] is False
    assert qa_data["sources"]
    assert "thinking_process" in qa_data

    favorite_resp = client.post(
        f"/api/v1/qa/{qa_data['record_id']}/favorite",
        json={"is_favorite": True},
        headers=student_headers,
    )
    assert favorite_resp.status_code == 200, favorite_resp.text
    feedback_resp = client.post(
        f"/api/v1/qa/{qa_data['record_id']}/feedback",
        json={"feedback": "positive", "feedback_comment": "解释清晰"},
        headers=student_headers,
    )
    assert feedback_resp.status_code == 200, feedback_resp.text
    history_resp = client.get(
        "/api/v1/qa/history",
        params={"course_id": course["id"], "lesson_id": lesson_id, "keyword": "矩阵"},
        headers=student_headers,
    )
    assert history_resp.status_code == 200, history_resp.text
    history = history_resp.json()["data"]
    assert len(history) >= 1
    assert history[0]["conversation_id"] == qa_data["conversation_id"]
    assert history[0]["record_count"] == 1
    assert "thinking_process" not in history[0]

    conversation_resp = client.get(f"/api/v1/qa/conversations/{qa_data['conversation_id']}", headers=student_headers)
    assert conversation_resp.status_code == 200, conversation_resp.text
    conversation_records = conversation_resp.json()["data"]
    assert len(conversation_records) == 1
    assert "thinking_process" in conversation_records[0]

    problem_resp = client.post(
        "/api/v1/tutoring/problems/text",
        json={"course_id": course["id"], "text": "已知矩阵A，求其行列式"},
        headers=student_headers,
    )
    assert problem_resp.status_code == 200, problem_resp.text
    problem = problem_resp.json()["data"]
    assert problem["knowledge_points"]

    guidance_lvl1 = client.get(
        f"/api/v1/tutoring/problems/{problem['id']}/guidance",
        params={"level": 1},
        headers=student_headers,
    )
    assert guidance_lvl1.status_code == 200, guidance_lvl1.text
    # #16：分层引导服务端顺序解锁，不能跳级，须逐级解锁 1 -> 2 -> 3。
    guidance_skip = client.get(
        f"/api/v1/tutoring/problems/{problem['id']}/guidance",
        params={"level": 3},
        headers=student_headers,
    )
    assert guidance_skip.status_code == 400, guidance_skip.text
    guidance_lvl2 = client.get(
        f"/api/v1/tutoring/problems/{problem['id']}/guidance",
        params={"level": 2},
        headers=student_headers,
    )
    assert guidance_lvl2.status_code == 200, guidance_lvl2.text
    guidance_lvl3 = client.get(
        f"/api/v1/tutoring/problems/{problem['id']}/guidance",
        params={"level": 3},
        headers=student_headers,
    )
    assert guidance_lvl3.status_code == 200, guidance_lvl3.text

    knowledge_resp = client.get(
        "/api/v1/learning/knowledge-points",
        params={"course_id": course["id"], "chapter_id": chapter["id"]},
        headers=student_headers,
    )
    assert knowledge_resp.status_code == 200, knowledge_resp.text
    assert len(knowledge_resp.json()["data"]) >= 1

    quiz_generate_resp = client.post(
        "/api/v1/learning/quizzes/generate",
        json={
            "course_id": course["id"],
            "chapter_id": chapter["id"],
            "title": "第一章测验",
            "quiz_type": "course",
            "question_count": 3,
        },
        headers=teacher_headers,
    )
    assert quiz_generate_resp.status_code == 200, quiz_generate_resp.text
    quiz = quiz_generate_resp.json()["data"]

    publish_quiz_resp = client.post(f"/api/v1/learning/quizzes/{quiz['id']}/publish", headers=teacher_headers)
    assert publish_quiz_resp.status_code == 200, publish_quiz_resp.text

    quiz_list_resp = client.get("/api/v1/learning/quizzes", params={"course_id": course["id"]}, headers=student_headers)
    assert quiz_list_resp.status_code == 200, quiz_list_resp.text
    assert len(quiz_list_resp.json()["data"]) >= 1

    quiz_detail_resp = client.get(f"/api/v1/learning/quizzes/{quiz['id']}", headers=student_headers)
    assert quiz_detail_resp.status_code == 200, quiz_detail_resp.text
    questions = quiz_detail_resp.json()["data"]["questions"]
    assert len(questions) == 3
    assert questions[0]["reference_answer"] is None

    submit_resp = client.post(
        f"/api/v1/learning/quizzes/{quiz['id']}/submit",
        json={
            "answers": [
                {"question_id": questions[0]["id"], "answer": 1},
                {"question_id": questions[1]["id"], "answer": 0},
                {"question_id": questions[2]["id"], "answer": "矩阵 线性 变换"},
            ]
        },
        headers=student_headers,
    )
    assert submit_resp.status_code == 200, submit_resp.text
    assert submit_resp.json()["data"]["score"] > 0

    wrong_resp = client.get("/api/v1/learning/wrong-questions", params={"course_id": course["id"]}, headers=student_headers)
    assert wrong_resp.status_code == 200, wrong_resp.text
    assert len(wrong_resp.json()["data"]) >= 1
    assert "knowledge_point_id" in wrong_resp.json()["data"][0]
    wrong_count_before = len(wrong_resp.json()["data"])

    wrong_practice_resp = client.post(
        "/api/v1/learning/wrong-questions/practice",
        params={"course_id": course["id"]},
        headers=student_headers,
    )
    assert wrong_practice_resp.status_code == 200, wrong_practice_resp.text
    wrong_practice = wrong_practice_resp.json()["data"]
    wrong_detail_resp = client.get(f"/api/v1/learning/quizzes/{wrong_practice['id']}", headers=student_headers)
    assert wrong_detail_resp.status_code == 200, wrong_detail_resp.text
    wrong_detail_questions = wrong_detail_resp.json()["data"]["questions"]
    assert wrong_detail_questions

    with db_session.SessionLocal() as db:
        answer_payload = []
        for item in wrong_detail_questions:
            stored_question = db.get(QuizQuestion, item["id"])
            reference = stored_question.reference_answer
            if isinstance(reference, dict):
                expected = reference.get("value", reference.get("answer", reference.get("correct_answer", reference.get("keywords", ""))))
            else:
                expected = reference
            if isinstance(expected, list):
                answer = " ".join(str(value) for value in expected) if item["question_type"] == "short_answer" else expected
            else:
                answer = expected
            answer_payload.append({"question_id": item["id"], "answer": answer})
    wrong_submit_resp = client.post(
        f"/api/v1/learning/quizzes/{wrong_practice['id']}/submit",
        json={"answers": answer_payload, "duration_seconds": 123},
        headers=student_headers,
    )
    assert wrong_submit_resp.status_code == 200, wrong_submit_resp.text
    assert wrong_submit_resp.json()["data"]["duration_seconds"] == 123
    wrong_after_resp = client.get("/api/v1/learning/wrong-questions", params={"course_id": course["id"]}, headers=student_headers)
    assert wrong_after_resp.status_code == 200, wrong_after_resp.text
    wrong_after_items = wrong_after_resp.json()["data"]
    assert len(wrong_after_items) == wrong_count_before
    # 掌握状态机：答对一次仅进入"巩固中"，不直接摘掉错题。
    assert not any(item["is_resolved"] for item in wrong_after_items)
    assert any(item["mastery"] == "consolidating" and item["correct_streak"] == 1 for item in wrong_after_items)
    assert all(item["history_count"] >= 1 for item in wrong_after_items)

    # 整卷重做：克隆出一份新私有练习，作答后连对次数累计到 2。艾宾浩斯曲线下"已掌握"需走完
    # 整条 [1,2,4,7,15,30] 天曲线，故此时仍为"巩固中"，并已排定下一档复习时间(next_review_at)。
    retake_resp = client.post(
        f"/api/v1/learning/quizzes/{wrong_practice['id']}/retake",
        json={"mode": "full"},
        headers=student_headers,
    )
    assert retake_resp.status_code == 200, retake_resp.text
    retake_quiz_data = retake_resp.json()["data"]
    assert retake_quiz_data["id"] != wrong_practice["id"]
    assert retake_quiz_data["question_count"] == len(wrong_detail_questions)
    retake_detail_resp = client.get(f"/api/v1/learning/quizzes/{retake_quiz_data['id']}", headers=student_headers)
    assert retake_detail_resp.status_code == 200, retake_detail_resp.text
    retake_questions = retake_detail_resp.json()["data"]["questions"]
    assert len(retake_questions) == len(wrong_detail_questions)
    # 学生视角作答前不泄露答案。
    assert all(item["reference_answer"] is None for item in retake_questions)
    with db_session.SessionLocal() as db:
        retake_answers = []
        for item in retake_questions:
            stored_question = db.get(QuizQuestion, item["id"])
            reference = stored_question.reference_answer
            if isinstance(reference, dict):
                expected = reference.get("value", reference.get("answer", reference.get("correct_answer", reference.get("keywords", ""))))
            else:
                expected = reference
            if isinstance(expected, list):
                answer = " ".join(str(value) for value in expected) if item["question_type"] == "short_answer" else expected
            else:
                answer = expected
            retake_answers.append({"question_id": item["id"], "answer": answer})
    retake_submit_resp = client.post(
        f"/api/v1/learning/quizzes/{retake_quiz_data['id']}/submit",
        json={"answers": retake_answers},
        headers=student_headers,
    )
    assert retake_submit_resp.status_code == 200, retake_submit_resp.text
    resolved_resp = client.get("/api/v1/learning/wrong-questions", params={"course_id": course["id"]}, headers=student_headers)
    assert resolved_resp.status_code == 200, resolved_resp.text
    resolved_items = resolved_resp.json()["data"]
    # 连对 2 次仍在曲线中途：巩固中、未掌握归档，且已排定下一档复习时间。
    assert any(
        (not item["is_resolved"])
        and item["mastery"] == "consolidating"
        and item["correct_streak"] == 2
        and item["review_stage"] == 2
        and item["next_review_at"]
        for item in resolved_items
    )

    # 生成任务占位卡列表：接口可用且只返回本人任务。
    tasks_resp = client.get(
        "/api/v1/learning/generation-tasks",
        params={"course_id": course["id"]},
        headers=student_headers,
    )
    assert tasks_resp.status_code == 200, tasks_resp.text
    assert isinstance(tasks_resp.json()["data"], list)

    # 删除练习：未开始的克隆卷可删；已作答的卷不可删。
    fresh_retake_resp = client.post(
        f"/api/v1/learning/quizzes/{wrong_practice['id']}/retake",
        json={"mode": "full"},
        headers=student_headers,
    )
    assert fresh_retake_resp.status_code == 200, fresh_retake_resp.text
    fresh_quiz_id = fresh_retake_resp.json()["data"]["id"]
    # 未开始 → 删除成功，且删除后取详情 404。
    del_resp = client.delete(f"/api/v1/learning/quizzes/{fresh_quiz_id}", headers=student_headers)
    assert del_resp.status_code == 200, del_resp.text
    assert client.get(f"/api/v1/learning/quizzes/{fresh_quiz_id}", headers=student_headers).status_code == 404
    # 已作答的错题重练卷 → 删除被拒。
    del_attempted_resp = client.delete(f"/api/v1/learning/quizzes/{wrong_practice['id']}", headers=student_headers)
    assert del_attempted_resp.status_code == 400, del_attempted_resp.text

    practice_resp = client.post(
        "/api/v1/learning/quizzes/generate",
        json={
            "course_id": course["id"],
            "chapter_id": chapter["id"],
            "chapter_ids": [chapter["id"]],
            "title": "章节自练",
            "quiz_type": "practice",
            "question_count": 2,
            "prefer_weak_points": True,
        },
        headers=student_headers,
    )
    assert practice_resp.status_code == 200, practice_resp.text
    practice_quiz = practice_resp.json()["data"]
    assert practice_quiz["task_id"]
    student_notifications_resp = client.get("/api/v1/student/notifications", headers=student_headers)
    assert student_notifications_resp.status_code == 200, student_notifications_resp.text
    assert any(
        item["type"] == "quiz_generated" and item.get("resource_id") == practice_quiz["id"]
        for item in student_notifications_resp.json()["data"]
    )

    weak_resp = client.get("/api/v1/learning/weak-points", params={"course_id": course["id"]}, headers=student_headers)
    assert weak_resp.status_code == 200, weak_resp.text
    assert isinstance(weak_resp.json()["data"], list)

    plan_resp = client.post(
        "/api/v1/learning/plans",
        json={
            "course_id": course["id"],
            "title": "期中复习计划",
            "goal": "一周内完成矩阵与行列式复习",
            "available_days": 3,
            "daily_minutes": 40,
        },
        headers=student_headers,
    )
    assert plan_resp.status_code == 200, plan_resp.text
    plan_data = plan_resp.json()["data"]
    assert len(plan_data["tasks"]) == 3
    first_task_id = plan_data["tasks"][0]["id"]

    checkin_resp = client.post(
        f"/api/v1/learning/tasks/{first_task_id}/checkin",
        json={"notes": "已完成第一天任务"},
        headers=student_headers,
    )
    assert checkin_resp.status_code == 200, checkin_resp.text

    records_resp = client.get("/api/v1/learning/records", params={"course_id": course["id"]}, headers=student_headers)
    assert records_resp.status_code == 200, records_resp.text
    records = records_resp.json()["data"]
    assert records["progress_count"] >= 1
    assert records["qa_count"] >= 1
    assert records["problem_count"] >= 1
    assert records["attempt_count"] >= 1
    assert records["recent_progress"]
    assert records["recent_qa"]
    assert records["recent_problems"]
    assert records["recent_attempts"]


def test_qa_concept_question_contributes_to_weak_points_and_quiz_priority(client, monkeypatch):
    course, chapter, lesson_id, _teacher_headers, student_headers = bootstrap_course_with_material(client)
    with db_session.SessionLocal() as db:
        db.execute(delete(KnowledgePoint).where(KnowledgePoint.course_id == course["id"]))
        db.add_all(
            [
                KnowledgePoint(
                    course_id=course["id"],
                    chapter_id=chapter["id"],
                    name="矩阵",
                    description="矩阵相关知识点",
                    content_by_level={},
                ),
                KnowledgePoint(
                    course_id=course["id"],
                    chapter_id=chapter["id"],
                    name="行列式",
                    description="行列式相关知识点",
                    content_by_level={},
                ),
            ]
        )
        db.commit()

    lesson_detail_resp = client.get(f"/api/v1/lessons/{lesson_id}", headers=student_headers)
    assert lesson_detail_resp.status_code == 200, lesson_detail_resp.text
    first_page_id = lesson_detail_resp.json()["data"]["pages"][0]["id"]

    qa_resp = client.post(
        "/api/v1/qa/ask",
        json={"course_id": course["id"], "lesson_page_id": first_page_id, "question": "行列式是什么概念"},
        headers=student_headers,
    )
    assert qa_resp.status_code == 200, qa_resp.text
    assert qa_resp.json()["data"]["is_out_of_scope"] is False

    with db_session.SessionLocal() as db:
        signal = db.scalar(
            select(StudentLearningSignal)
            .join(KnowledgePoint, KnowledgePoint.id == StudentLearningSignal.knowledge_point_id)
            .where(StudentLearningSignal.course_id == course["id"], KnowledgePoint.name == "行列式")
        )
        assert signal is not None
        assert signal.intent == "concept_confusion"
        assert signal.score > 0

    weak_resp = client.get("/api/v1/learning/weak-points", params={"course_id": course["id"]}, headers=student_headers)
    assert weak_resp.status_code == 200, weak_resp.text
    weak_items = weak_resp.json()["data"]
    determinant = next(item for item in weak_items if item["knowledge_point"] == "行列式")
    assert determinant["wrong_count"] == 0
    assert determinant["qa_signal_count"] >= 1
    assert determinant["weak_score"] > 0

    captured = {}

    def fake_prioritized_quiz_questions(*, topic, source_text, count, db=None, **kwargs):
        captured["topic"] = topic
        return fake_quiz_questions(topic=topic, source_text=source_text, count=count, db=db)

    monkeypatch.setattr(ai_service, "generate_quiz_questions", fake_prioritized_quiz_questions)
    quiz_resp = client.post(
        "/api/v1/learning/quizzes/generate",
        json={
            "course_id": course["id"],
            "chapter_id": chapter["id"],
            "chapter_ids": [chapter["id"]],
            "title": "薄弱点章节练习",
            "quiz_type": "practice",
            "question_count": 1,
            "prefer_weak_points": True,
        },
        headers=student_headers,
    )
    assert quiz_resp.status_code == 200, quiz_resp.text
    assert captured["topic"].startswith("行列式")


def test_teacher_weak_quiz_management_flow(client, monkeypatch):
    monkeypatch.setattr(ai_service, "generate_quiz_questions", fake_quiz_questions_with_types)
    course, chapter, _lesson_id, teacher_headers, student_headers = bootstrap_course_with_material(client)

    base_quiz_resp = client.post(
        "/api/v1/learning/quizzes/generate",
        json={
            "course_id": course["id"],
            "chapter_id": chapter["id"],
            "title": "薄弱点来源测验",
            "quiz_type": "course",
            "question_count": 1,
        },
        headers=teacher_headers,
    )
    assert base_quiz_resp.status_code == 200, base_quiz_resp.text
    base_quiz = base_quiz_resp.json()["data"]
    publish_resp = client.post(f"/api/v1/learning/quizzes/{base_quiz['id']}/publish", headers=teacher_headers)
    assert publish_resp.status_code == 200, publish_resp.text
    detail_resp = client.get(f"/api/v1/learning/quizzes/{base_quiz['id']}", headers=student_headers)
    assert detail_resp.status_code == 200, detail_resp.text
    source_question = detail_resp.json()["data"]["questions"][0]
    submit_resp = client.post(
        f"/api/v1/learning/quizzes/{base_quiz['id']}/submit",
        json={"answers": [{"question_id": source_question["id"], "answer": 1}]},
        headers=student_headers,
    )
    assert submit_resp.status_code == 200, submit_resp.text

    weak_list_resp = client.get("/api/v1/learning/teacher/weak-quizzes", params={"course_id": course["id"]}, headers=teacher_headers)
    assert weak_list_resp.status_code == 200, weak_list_resp.text
    weak_payload = weak_list_resp.json()["data"]
    assert weak_payload["stats"]["weak_point_count"] >= 1
    weak_point = weak_payload["weak_points"][0]

    generate_resp = client.post(
        "/api/v1/learning/teacher/weak-quizzes/generate",
        json={
            "course_id": course["id"],
            "weak_point_id": weak_point["knowledge_point_id"],
            "title": "矩阵薄弱点专项",
            "question_count": 2,
            "question_type_counts": {"single_choice": 1, "judge": 1},
        },
        headers=teacher_headers,
    )
    assert generate_resp.status_code == 200, generate_resp.text
    weak_quiz = generate_resp.json()["data"]
    assert weak_quiz["task_id"]
    assert weak_quiz["metadata_json"]["weak_quiz"] is True
    assert weak_quiz["metadata_json"]["question_type_counts"] == {"single_choice": 1, "judge": 1}
    generation_task_resp = client.get(
        f"/api/v1/learning/generation-tasks/{weak_quiz['task_id']}",
        headers=teacher_headers,
    )
    assert generation_task_resp.status_code == 200, generation_task_resp.text
    assert generation_task_resp.json()["data"]["id"] == weak_quiz["id"]
    teacher_dashboard_resp = client.get("/api/v1/teacher/dashboard", headers=teacher_headers)
    assert teacher_dashboard_resp.status_code == 200, teacher_dashboard_resp.text
    assert any(
        item["type"] == "quiz_generated" and item.get("resource_id") == weak_quiz["id"]
        for item in teacher_dashboard_resp.json()["data"]["notifications"]
    )

    attempts_empty_resp = client.get(f"/api/v1/learning/teacher/weak-quizzes/{weak_quiz['id']}/attempts", headers=teacher_headers)
    assert attempts_empty_resp.status_code == 200, attempts_empty_resp.text
    assert len(attempts_empty_resp.json()["data"]["questions"]) == 2

    publish_weak_resp = client.post(f"/api/v1/learning/quizzes/{weak_quiz['id']}/publish", headers=teacher_headers)
    assert publish_weak_resp.status_code == 200, publish_weak_resp.text
    weak_detail_resp = client.get(f"/api/v1/learning/quizzes/{weak_quiz['id']}", headers=student_headers)
    assert weak_detail_resp.status_code == 200, weak_detail_resp.text
    weak_questions = weak_detail_resp.json()["data"]["questions"]
    weak_submit_resp = client.post(
        f"/api/v1/learning/quizzes/{weak_quiz['id']}/submit",
        json={"answers": [{"question_id": item["id"], "answer": 0} for item in weak_questions]},
        headers=student_headers,
    )
    assert weak_submit_resp.status_code == 200, weak_submit_resp.text

    attempts_resp = client.get(f"/api/v1/learning/teacher/weak-quizzes/{weak_quiz['id']}/attempts", headers=teacher_headers)
    assert attempts_resp.status_code == 200, attempts_resp.text
    attempts_payload = attempts_resp.json()["data"]
    assert attempts_payload["quiz"]["attempt_count"] == 1
    assert attempts_payload["attempts"][0]["correct_count"] == 2


def test_quiz_generation_requires_ai_questions(monkeypatch):
    monkeypatch.setattr(ai_service, "_call_chat", lambda *args, **kwargs: None)
    clean = sanitize_quiz_source_text(
        "![d03903cfb024ca7b18336cdcfff771e3.jpeg](https://classagent.oss-cn-beijing.aliyuncs.com/a.png)\n"
        "## 语法制导翻译概述\n语义分析的任务包括语义检查和翻译。"
    )
    assert "https" not in clean
    assert "d03903cfb024" not in clean
    assert "jpeg" not in clean
    assert "语义分析" in clean
    with pytest.raises(AppError) as exc:
        ai_service.generate_quiz_questions(topic="语法制导翻译", source_text=clean, count=1)
    assert "AI 出题失败" in exc.value.detail["message"]


def test_quiz_generation_normalizes_ai_reference_answer(monkeypatch):
    def ai_content(*args, **kwargs):
        return (
            '{"items":[{"question_type":"single_choice",'
            '"stem":"属性文法三元组 A=(G,C,F) 中，F 表示什么？",'
            '"options":["上下文无关文法","属性的计算规则","属性的有穷集","词法分析规则"],'
            '"reference_answer":{"correct_option":"B"},'
            '"explanation":"资料中说明 F 是关于属性的计算规则。",'
            '"score":10,"difficulty":"standard"}]}'
        )

    monkeypatch.setattr(ai_service, "_call_chat", ai_content)
    questions = ai_service.generate_quiz_questions(
        topic="属性文法",
        source_text="属性文法定义为 A=(G,C,F)，其中 F 是关于属性的计算规则。",
        count=1,
    )
    assert questions[0]["reference_answer"] == {"value": 1}


def test_quiz_generation_accepts_chinese_ai_field_aliases(monkeypatch):
    def ai_content(*args, **kwargs):
        return (
            '{"items":[{"题型":"单选题",'
            '"题干":"属性文法三元组 A=(G,C,F) 中，F 表示什么？",'
            '"选项":["上下文无关文法","属性的计算规则","属性的有穷集","词法分析规则"],'
            '"正确答案":"B",'
            '"解析":"资料中说明 F 是关于属性的计算规则。",'
            '"score":10,"difficulty":"standard"}]}'
        )

    monkeypatch.setattr(ai_service, "_call_chat", ai_content)
    questions = ai_service.generate_quiz_questions(
        topic="属性文法",
        source_text="属性文法定义为 A=(G,C,F)，其中 F 是关于属性的计算规则。",
        count=1,
    )
    assert questions[0]["question_type"] == "single_choice"
    assert questions[0]["reference_answer"] == {"value": 1}
    assert questions[0]["explanation"] == "资料中说明 F 是关于属性的计算规则。"


def test_quiz_generation_retries_ai_when_valid_questions_are_insufficient(monkeypatch):
    calls = iter(
        [
            (
                '{"items":[{"question_type":"short_answer",'
                '"stem":"在自底向上的语法分析中，语义动作在何时执行？",'
                '"options":null,'
                '"reference_answer":{"keywords":["归约","产生式","语义动作"]},'
                '"explanation":"资料中说明自底向上分析在规约时执行语义动作。",'
                '"score":10,"difficulty":"standard"}]}'
            ),
            (
                '{"items":[{"question_type":"single_choice",'
                '"stem":"在自底向上的语法分析中，语义动作在什么时机执行？",'
                '"options":["进行派生时","进行归约时","语法分析开始前","语法分析结束后"],'
                '"reference_answer":{"value":1},'
                '"explanation":"资料中说明自底向上分析在归约时执行语义动作。",'
                '"score":10,"difficulty":"standard"}]}'
            ),
        ]
    )

    monkeypatch.setattr(ai_service, "_call_chat", lambda *args, **kwargs: next(calls))
    questions = ai_service.generate_quiz_questions(
        topic="语法制导翻译",
        source_text="自底向上的分析中，语义动作和产生式关联，当用产生式进行归约时执行。",
        count=1,
    )
    assert questions[0]["question_type"] == "single_choice"
    assert questions[0]["reference_answer"] == {"value": 1}


def test_quiz_generation_practice_fast_cuts_llm_calls(monkeypatch):
    # 学生自助练习提速档应把最坏 3 次串行大模型调用（主生成 + critic 自评 + 定向重试）降到 1 次。
    def batch(n=14):
        items = [
            {
                "question_type": "single_choice",
                "stem": f"下列关于编译原理概念 C{i} 的表述中，正确的是哪一项？",
                "options": [
                    f"C{i} 的正确定义表述",
                    f"C{i} 的错误表述甲项",
                    f"C{i} 的错误表述乙项",
                    f"C{i} 的错误表述丙项",
                ],
                "reference_answer": {"value": 0},
                "explanation": f"正确项准确刻画了 C{i}，其余三项分别在范围、条件、对象上出现常见混淆。",
                "difficulty": "standard",
                "knowledge_point": f"概念C{i}",
                "cognitive_level": "理解",
            }
            for i in range(n)
        ]
        return json.dumps({"items": items}, ensure_ascii=False)

    calls: list = []

    def fake_call_chat(*args, **kwargs):
        calls.append(kwargs.get("purpose"))
        return batch()

    monkeypatch.setattr(ai_service, "_call_chat", fake_call_chat)
    source = "属性文法定义为 A=(G,C,F)，其中 F 是关于属性的计算规则。语义分析包括语义检查与翻译。"

    calls.clear()
    fast = ai_service.generate_quiz_questions(topic="编译原理", source_text=source, count=5, practice_fast=True)
    fast_calls = len(calls)

    calls.clear()
    normal = ai_service.generate_quiz_questions(topic="编译原理", source_text=source, count=5, practice_fast=False)
    normal_calls = len(calls)

    assert len(fast) == 5
    assert len(normal) == 5
    # 提速档只调用一次主生成；默认档额外触发 critic 自评（且通常再加一次定向重试）
    assert fast_calls == 1
    assert normal_calls > fast_calls


def test_quiz_generation_respects_question_type_counts(monkeypatch):
    def ai_content(*args, **kwargs):
        return (
            '{"items":['
            '{"question_type":"single_choice",'
            '"stem":"矩阵可以表示线性变换，下列说法正确的是？",'
            '"options":["矩阵可以表示线性变换","矩阵只能表示图片","矩阵与变换无关","矩阵是文件名"],'
            '"reference_answer":{"value":0},'
            '"explanation":"资料中说明矩阵可以表示线性变换。",'
            '"score":10,"difficulty":"standard"},'
            '{"question_type":"blank",'
            '"stem":"矩阵可以表示____。",'
            '"options":null,'
            '"reference_answer":{"keywords":["线性变换"]},'
            '"explanation":"填入线性变换。",'
            '"score":10,"difficulty":"standard"}]}'
        )

    monkeypatch.setattr(ai_service, "_call_chat", ai_content)
    questions = ai_service.generate_quiz_questions(
        topic="矩阵",
        source_text="矩阵可以表示线性变换，行列式反映缩放系数。",
        count=2,
        type_counts={"single_choice": 1, "blank": 1},
    )
    assert [item["question_type"] for item in questions] == ["single_choice", "blank"]
    assert questions[1]["reference_answer"] == {"keywords": ["线性变换"]}


def test_quiz_generation_prioritizes_applied_questions(monkeypatch):
    captured = {}

    def ai_content(*args, **kwargs):
        captured["prompt"] = kwargs["user_prompt"]
        return (
            '{"items":['
            '{"question_type":"single_choice",'
            '"stem":"矩阵可以表示什么？",'
            '"options":["线性变换","数字表格","坐标变化","向量集合"],'
            '"reference_answer":{"value":0},'
            '"explanation":"资料说明矩阵可以表示线性变换。",'
            '"score":10,"difficulty":"standard"},'
            '{"question_type":"single_choice",'
            '"stem":"已知线性变换把向量(1,0)映射到(2,0)，把(0,1)映射到(0,3)，应如何用矩阵表示并判断面积缩放？",'
            '"options":["矩阵为[[2,0],[0,3]]，面积放大6倍","矩阵为[[3,0],[0,2]]，面积放大5倍","矩阵为[[2,3],[0,0]]，面积不变","矩阵为[[0,2],[3,0]]，面积缩小6倍"],'
            '"reference_answer":{"value":0},'
            '"explanation":"矩阵的列向量对应基向量像，行列式为6。",'
            '"score":10,"difficulty":"standard"}]}'
        )

    monkeypatch.setattr(ai_service, "_call_chat", ai_content)
    questions = ai_service.generate_quiz_questions(
        topic="矩阵与行列式",
        source_text="矩阵可以表示线性变换，行列式反映面积缩放系数。公式 det(A)=ad-bc。",
        count=1,
    )
    assert "应用题" in captured["prompt"]
    assert "纯概念题最多" in captured["prompt"]
    assert "计算题" in captured["prompt"]
    assert questions[0]["stem"].startswith("已知线性变换")


def test_quiz_generation_rejects_direct_fact_short_answer(monkeypatch):
    def ai_content(*args, **kwargs):
        return (
            '{"items":[{"question_type":"short_answer",'
            '"stem":"在自底向上的语法分析中，语义动作在何时执行？",'
            '"options":null,'
            '"reference_answer":{"keywords":["归约","产生式","语义动作"]},'
            '"explanation":"资料中说明自底向上分析在规约时执行语义动作。",'
            '"score":10,"difficulty":"standard"}]}'
        )

    monkeypatch.setattr(ai_service, "_call_chat", ai_content)
    with pytest.raises(AppError) as exc:
        ai_service.generate_quiz_questions(
            topic="语法制导翻译",
            source_text="自底向上的分析中，语义动作和产生式关联，当用产生式进行归约时执行。",
            count=1,
        )
    assert "有效题目不足" in exc.value.detail["message"]


def test_qa_stream_falls_back_to_regular_answer_when_upstream_stream_fails(client, monkeypatch):
    course, _chapter, lesson_id, _teacher_headers, student_headers = bootstrap_course_with_material(client)

    from app.services import qa as qa_service

    lesson_detail_resp = client.get(f"/api/v1/lessons/{lesson_id}", headers=student_headers)
    assert lesson_detail_resp.status_code == 200, lesson_detail_resp.text
    first_page_id = lesson_detail_resp.json()["data"]["pages"][0]["id"]

    def fail_stream(**kwargs):
        raise RuntimeError("stream unavailable")

    def fallback_answer(**kwargs):
        return "非流式回退回答：矩阵可以表示线性变换。这个结论说明矩阵不仅是数字表格，还可以描述向量在空间中的旋转、缩放和投影。", False, "非流式回退思考过程"

    monkeypatch.setattr(qa_service.ai_service, "stream_answer_question", fail_stream)
    monkeypatch.setattr(qa_service.ai_service, "answer_question", fallback_answer)

    with client.stream(
        "POST",
        "/api/v1/qa/ask/stream",
        json={"course_id": course["id"], "lesson_page_id": first_page_id, "question": "矩阵可以表示什么"},
        headers=student_headers,
    ) as response:
        body = "".join(response.iter_text())

    assert response.status_code == 200, body
    assert "event: error" not in body
    assert "非流式回退回答" in body
    assert "非流式回退思考过程" in body
    assert "event: final" in body
    assert body.count("event: delta") >= 2


def test_qa_stream_final_does_not_wait_for_learning_signal_postprocess(client, monkeypatch):
    import threading
    from time import perf_counter

    course, _chapter, _lesson_id, _teacher_headers, student_headers = bootstrap_course_with_material(client)

    from app.services import qa as qa_service
    from app.services.ai import ChatDelta

    def stream_answer(**kwargs):
        yield ChatDelta("content", "行列式反映线性变换中的缩放系数。")

    entered = threading.Event()
    release = threading.Event()

    def slow_record_signals(db, *, user, record):
        entered.set()
        release.wait(timeout=5)
        return []

    monkeypatch.setattr(qa_service.ai_service, "stream_answer_question", stream_answer)
    monkeypatch.setattr(qa_service, "record_qa_learning_signals", slow_record_signals)

    start = perf_counter()
    try:
        with client.stream(
            "POST",
            "/api/v1/qa/ask/stream",
            json={"course_id": course["id"], "question": "行列式是什么概念？"},
            headers=student_headers,
        ) as response:
            body = "".join(response.iter_text())
        elapsed = perf_counter() - start

        assert response.status_code == 200, body
        assert "event: final" in body
        assert "行列式反映线性变换中的缩放系数" in body
        assert entered.wait(timeout=1)
        assert elapsed < 2
    finally:
        release.set()


def test_course_qa_history_excludes_lesson_page_conversations_by_default(client, monkeypatch):
    course, _chapter, lesson_id, _teacher_headers, student_headers = bootstrap_course_with_material(client)

    from app.services import qa as qa_service

    lesson_detail_resp = client.get(f"/api/v1/lessons/{lesson_id}", headers=student_headers)
    assert lesson_detail_resp.status_code == 200, lesson_detail_resp.text
    first_page_id = lesson_detail_resp.json()["data"]["pages"][0]["id"]

    monkeypatch.setattr(
        qa_service.ai_service,
        "answer_question",
        lambda **kwargs: ("课件问答回答：矩阵可以表示线性变换。", False, None),
    )
    ask_resp = client.post(
        "/api/v1/qa/ask",
        json={"course_id": course["id"], "lesson_page_id": first_page_id, "question": "课件里的矩阵是什么？"},
        headers=student_headers,
    )
    assert ask_resp.status_code == 200, ask_resp.text
    lesson_conversation_id = ask_resp.json()["data"]["conversation_id"]

    mixed_scope_resp = client.post(
        "/api/v1/qa/ask",
        json={"course_id": course["id"], "conversation_id": lesson_conversation_id, "question": "同一课件对话后续追问"},
        headers=student_headers,
    )
    assert mixed_scope_resp.status_code == 200, mixed_scope_resp.text

    history_resp = client.get("/api/v1/qa/history", params={"course_id": course["id"]}, headers=student_headers)
    assert history_resp.status_code == 200, history_resp.text
    history = history_resp.json()["data"]
    assert all(item["lesson_page_id"] is None for item in history)
    assert not any(item["question"] == "课件里的矩阵是什么？" for item in history)
    assert not any(item["question"] == "同一课件对话后续追问" for item in history)

    lesson_history_resp = client.get(
        "/api/v1/qa/history",
        params={"course_id": course["id"], "lesson_id": lesson_id},
        headers=student_headers,
    )
    assert lesson_history_resp.status_code == 200, lesson_history_resp.text
    lesson_history = lesson_history_resp.json()["data"]
    assert [item["conversation_id"] for item in lesson_history] == [lesson_conversation_id]
    assert lesson_history[0]["question"] == "课件里的矩阵是什么？"
    assert lesson_history[0]["lesson_page_id"] == first_page_id
    assert lesson_history[0]["record_count"] == 2

    lesson_conversation_resp = client.get(f"/api/v1/qa/conversations/{lesson_conversation_id}", headers=student_headers)
    assert lesson_conversation_resp.status_code == 200, lesson_conversation_resp.text
    assert [item["question"] for item in lesson_conversation_resp.json()["data"]] == ["课件里的矩阵是什么？", "同一课件对话后续追问"]


def test_qa_stream_sends_previous_turn_as_history_messages(client, monkeypatch):
    course, _chapter, _lesson_id, _teacher_headers, student_headers = bootstrap_course_with_material(client)

    from app.services import qa as qa_service
    from app.services.ai import ChatDelta

    calls: list[dict] = []

    def stream_answer(**kwargs):
        calls.append(kwargs)
        if len(calls) == 1:
            yield ChatDelta("content", "矩阵可以表示线性变换。")
        else:
            yield ChatDelta("content", "这里的“这个”指上一轮提到的线性变换。")

    monkeypatch.setattr(qa_service.ai_service, "stream_answer_question", stream_answer)

    with client.stream(
        "POST",
        "/api/v1/qa/ask/stream",
        json={"course_id": course["id"], "question": "矩阵可以表示什么？"},
        headers=student_headers,
    ) as response:
        first_body = "".join(response.iter_text())

    assert response.status_code == 200, first_body
    assert "event: final" in first_body
    final_payload = first_body.split("event: final\ndata: ", 1)[1].split("\n\n", 1)[0]
    conversation_id = json.loads(final_payload)["conversation_id"]

    with client.stream(
        "POST",
        "/api/v1/qa/ask/stream",
        json={"course_id": course["id"], "conversation_id": conversation_id, "question": "这个能再举个例子吗？"},
        headers=student_headers,
    ) as response:
        second_body = "".join(response.iter_text())

    assert response.status_code == 200, second_body
    assert len(calls) == 2
    assert calls[0]["history"] == []
    assert calls[1]["history"] == [
        {"role": "user", "content": "矩阵可以表示什么？"},
        {"role": "assistant", "content": "矩阵可以表示线性变换。"},
    ]
    assert calls[1]["question"] == "这个能再举个例子吗？"
    assert "线性变换" in second_body

    history_resp = client.get("/api/v1/qa/history", params={"course_id": course["id"]}, headers=student_headers)
    assert history_resp.status_code == 200, history_resp.text
    history = history_resp.json()["data"]
    assert [item["conversation_id"] for item in history] == [conversation_id]
    assert history[0]["record_count"] == 2

    records_resp = client.get(f"/api/v1/qa/conversations/{conversation_id}", headers=student_headers)
    assert records_resp.status_code == 200, records_resp.text
    assert [item["question"] for item in records_resp.json()["data"]] == ["矩阵可以表示什么？", "这个能再举个例子吗？"]


def test_qa_from_lesson_page_can_answer_other_page_in_same_ppt(client, monkeypatch):
    course, _chapter, lesson_id, _teacher_headers, student_headers = bootstrap_course_with_material(client)

    from app.services import qa as qa_service
    from app.services.ai import ChatDelta

    lesson_detail_resp = client.get(f"/api/v1/lessons/{lesson_id}", headers=student_headers)
    assert lesson_detail_resp.status_code == 200, lesson_detail_resp.text
    first_page_id = lesson_detail_resp.json()["data"]["pages"][0]["id"]

    captured: dict[str, list[str]] = {}
    monkeypatch.setattr(qa_service, "search_course_knowledge", lambda *args, **kwargs: [])

    def stream_answer(**kwargs):
        captured["contexts"] = list(kwargs["contexts"])
        yield ChatDelta("content", "第2页讲行列式反映缩放系数。")

    monkeypatch.setattr(qa_service.ai_service, "stream_answer_question", stream_answer)

    with client.stream(
        "POST",
        "/api/v1/qa/ask/stream",
        json={"course_id": course["id"], "lesson_page_id": first_page_id, "question": "第2页讲了什么？"},
        headers=student_headers,
    ) as response:
        body = "".join(response.iter_text())

    assert response.status_code == 200, body
    assert "行列式反映缩放系数" in body
    assert any("第2页" in item and "行列式反映缩放系数" in item for item in captured["contexts"])


def test_qa_from_lesson_page_retrieves_with_current_lesson_scope(client, monkeypatch):
    course, _chapter, lesson_id, _teacher_headers, student_headers = bootstrap_course_with_material(client)

    from app.services import qa as qa_service

    lesson_detail_resp = client.get(f"/api/v1/lessons/{lesson_id}", headers=student_headers)
    assert lesson_detail_resp.status_code == 200, lesson_detail_resp.text
    first_page_id = lesson_detail_resp.json()["data"]["pages"][0]["id"]

    calls: list[dict] = []

    def fake_search(*args, **kwargs):
        calls.append(kwargs)
        return []

    monkeypatch.setattr(qa_service, "search_course_knowledge", fake_search)

    def answer_question(**kwargs):
        return "已根据当前课件上下文回答。", False, None

    monkeypatch.setattr(qa_service.ai_service, "answer_question", answer_question)
    response = client.post(
        "/api/v1/qa/ask",
        json={"course_id": course["id"], "lesson_page_id": first_page_id, "question": "第2页讲了什么？"},
        headers=student_headers,
    )

    assert response.status_code == 200, response.text
    assert calls
    assert calls[0]["lesson_id"] == lesson_id
    assert calls[0]["lesson_page_id"] is None
    assert calls[0]["limit"] >= 8


def test_qa_page_can_fallback_to_lesson_pages_without_page_context(client, monkeypatch):
    course, _chapter, _lesson_id, _teacher_headers, student_headers = bootstrap_course_with_material(client)

    from app.services import qa as qa_service

    captured: dict[str, list[str]] = {}
    monkeypatch.setattr(qa_service, "search_course_knowledge", lambda *args, **kwargs: [])

    def answer_question(**kwargs):
        captured["contexts"] = list(kwargs["contexts"])
        return "行列式反映缩放系数。", False, None

    monkeypatch.setattr(qa_service.ai_service, "answer_question", answer_question)

    response = client.post(
        "/api/v1/qa/ask",
        json={"course_id": course["id"], "question": "行列式反映什么？"},
        headers=student_headers,
    )

    assert response.status_code == 200, response.text
    data = response.json()["data"]
    assert data["is_out_of_scope"] is False
    assert "行列式反映缩放系数" in data["answer"]
    assert any("行列式反映缩放系数" in item for item in captured["contexts"])


def test_qa_uses_chapter_context_for_chapter_overview_when_vector_search_misses(client, monkeypatch):
    course, chapter, _lesson_id, _teacher_headers, student_headers = bootstrap_course_with_material(client)

    from app.services import qa as qa_service

    captured: dict[str, list[str]] = {}
    monkeypatch.setattr(qa_service, "search_course_knowledge", lambda *args, **kwargs: [])
    monkeypatch.setattr(
        qa_service.ai_service,
        "plan_qa_task",
        lambda **kwargs: {
            "scope": "chapter_overview",
            "question_type": "chapter_overview",
            "chapter_ids": [chapter["id"]],
            "chapter_id": chapter["id"],
            "keywords": ["矩阵基础"],
            "search_phrases": ["矩阵基础 重点"],
            "retrieval_query": "矩阵基础 重点",
            "tools": ["get_chapter_summary", "quote_source"],
        },
    )

    def answer_question(**kwargs):
        captured["contexts"] = list(kwargs["contexts"])
        return "已根据章节资料生成重点。", False, None

    monkeypatch.setattr(qa_service.ai_service, "answer_question", answer_question)
    response = client.post(
        "/api/v1/qa/ask",
        json={"course_id": course["id"], "question": "第一章 的重点是什么？"},
        headers=student_headers,
    )
    assert response.status_code == 200, response.text
    data = response.json()["data"]
    assert data["is_out_of_scope"] is False
    assert data["sources"]
    assert captured["contexts"]
    assert any(item.startswith("章节：") for item in captured["contexts"])
    assert any("资料片段：" in item or "页面内容：" in item or "资料：" in item for item in captured["contexts"])


def test_qa_chapter_range_uses_all_requested_chapters_when_vector_search_misses(client, monkeypatch):
    course, _chapter, _lesson_id, _teacher_headers, student_headers = bootstrap_course_with_material(client)

    with db_session.SessionLocal() as db:
        db_course = db.get(Course, course["id"])
        assert db_course is not None
        range_page_id = None
        for number in range(5, 8):
            chapter = Chapter(
                course_id=db_course.id,
                title=f"第{number}章 范围章节{number}",
                description=f"第{number}章说明",
                order_index=number,
            )
            db.add(chapter)
            db.flush()
            material = CourseMaterial(
                course_id=db_course.id,
                chapter_id=chapter.id,
                uploader_id=db_course.teacher_id,
                title=f"第{number}章课件",
                category=MaterialCategory.COURSEWARE.value,
                material_type=MaterialType.PDF.value,
                size_bytes=128,
                original_filename=f"chapter-{number}.pdf",
                storage_path=f"tests/chapter-{number}.pdf",
                extracted_text=f"第{number}章完整资料讲解重点{number}",
                parse_status=ProcessStatus.READY.value,
                vector_status=ProcessStatus.READY.value,
            )
            db.add(material)
            db.flush()
            lesson = Lesson(
                course_id=db_course.id,
                chapter_id=chapter.id,
                material_id=material.id,
                title=f"第{number}章课时",
                summary=f"第{number}章总结",
                page_count=2,
                status=LessonStatus.PUBLISHED.value,
            )
            db.add(lesson)
            db.flush()
            concept_page = LessonPage(
                lesson_id=lesson.id,
                page_number=1,
                page_title=f"第{number}章概念",
                page_text=f"第{number}章核心概念完整页面内容",
                script_text=f"第{number}章核心概念讲解文稿",
                script_status=ProcessStatus.READY.value,
            )
            method_page = LessonPage(
                lesson_id=lesson.id,
                page_number=2,
                page_title=f"第{number}章方法",
                page_text=f"第{number}章方法步骤完整页面内容",
                script_text=f"第{number}章方法步骤讲解文稿",
                script_status=ProcessStatus.READY.value,
            )
            db.add_all(
                [
                    concept_page,
                    method_page,
                ]
            )
            db.flush()
            if number == 5:
                range_page_id = concept_page.id
        db.commit()

    from app.services import qa as qa_service

    captured: dict[str, list[str]] = {}
    monkeypatch.setattr(qa_service, "search_course_knowledge", lambda *args, **kwargs: [])
    monkeypatch.setattr(
        qa_service.ai_service,
        "plan_qa_task",
        lambda **kwargs: {
            "scope": "chapter_overview",
            "question_type": "large_chapter_request",
            "chapter_ids": [chapter["id"] for chapter in kwargs["chapters"] if chapter["order_index"] in {5, 6, 7}],
            "keywords": ["范围章节"],
            "search_phrases": ["第5章 第6章 第7章"],
            "retrieval_query": "第5章 第6章 第7章",
            "tools": ["get_chapter_summary", "get_section_summary", "quote_source"],
            "large_request": True,
        },
    )

    def answer_question(**kwargs):
        captured["contexts"] = list(kwargs["contexts"])
        return "已讲解第5到第7章。", False, None

    monkeypatch.setattr(qa_service.ai_service, "answer_question", answer_question)
    assert range_page_id is not None
    response = client.post(
        "/api/v1/qa/ask",
        json={"course_id": course["id"], "lesson_page_id": range_page_id, "question": "请为我详细讲解5-7章的知识"},
        headers=student_headers,
    )
    assert response.status_code == 200, response.text
    joined_context = "\n".join(captured["contexts"])
    assert "第5章核心概念完整页面内容" in joined_context
    assert "第6章核心概念完整页面内容" in joined_context
    assert "第7章核心概念完整页面内容" in joined_context

def test_rag_context_packing_keeps_later_chapters_when_context_is_long():
    contexts = [f"第{number}章 " + (f"核心内容{number}" * 1200) for number in range(5, 8)]

    packed = _pack_rag_contexts(contexts, limit=1200)

    assert len(packed) <= 1200
    assert "第5章" in packed
    assert "第6章" in packed
    assert "第7章" in packed


def test_qa_uses_course_context_for_course_overview_when_vector_search_misses(client, monkeypatch):
    course, _chapter, _lesson_id, _teacher_headers, student_headers = bootstrap_course_with_material(client)

    from app.services import qa as qa_service

    captured: dict[str, list[str]] = {}
    monkeypatch.setattr(qa_service, "search_course_knowledge", lambda *args, **kwargs: [])
    monkeypatch.setattr(
        qa_service.ai_service,
        "plan_qa_task",
        lambda **kwargs: {
            "scope": "course_overview",
            "question_type": "course_overview",
            "chapter_ids": [],
            "chapter_id": None,
            "keywords": ["课程重点"],
            "search_phrases": ["课程重点 复习"],
            "retrieval_query": "课程重点 复习",
            "tools": ["get_chapter_summary", "get_section_summary", "quote_source"],
        },
    )

    def answer_question(**kwargs):
        captured["contexts"] = list(kwargs["contexts"])
        return "已根据课程资料生成整体重点。", False, None

    monkeypatch.setattr(qa_service.ai_service, "answer_question", answer_question)
    response = client.post(
        "/api/v1/qa/ask",
        json={"course_id": course["id"], "question": "这门课的重点是什么？"},
        headers=student_headers,
    )
    assert response.status_code == 200, response.text
    data = response.json()["data"]
    assert data["is_out_of_scope"] is False
    assert data["sources"]
    assert captured["contexts"]
    assert any(item.startswith("课程：") for item in captured["contexts"])
    assert any("章节结构：" in item or "资料片段：" in item or "页面内容：" in item or "资料：" in item for item in captured["contexts"])


def test_qa_course_review_question_uses_task_plan_and_reaches_final_answer(client, monkeypatch):
    course, _chapter, _lesson_id, _teacher_headers, student_headers = bootstrap_course_with_material(client)

    from app.services import qa as qa_service

    with db_session.SessionLocal() as db:
        db_course = db.get(Course, course["id"])
        assert db_course is not None
        db_course.name = "编译原理"
        db_course.description = "词法分析、语法分析、语义分析与代码生成"
        db.add(db_course)
        db.commit()

    calls = {"plan": 0, "answer": 0}
    captured: dict[str, list[str]] = {}
    monkeypatch.setattr(qa_service, "search_course_knowledge", lambda *args, **kwargs: [])

    def plan_qa_task(**kwargs):
        calls["plan"] += 1
        assert kwargs["question"] == "我该如何复习《编译原理》？"
        return {
            "scope": "course_overview",
            "question_type": "course_overview",
            "chapter_ids": [],
            "chapter_id": None,
            "keywords": ["编译原理", "复习"],
            "search_phrases": ["编译原理 复习"],
            "retrieval_query": "编译原理 复习提纲",
            "tools": ["get_chapter_summary", "get_section_summary", "quote_source"],
        }

    def answer_question(**kwargs):
        calls["answer"] += 1
        captured["contexts"] = list(kwargs["contexts"])
        return "建议按词法分析、语法分析、语义分析和代码生成建立复习框架。", False, None

    monkeypatch.setattr(qa_service.ai_service, "plan_qa_task", plan_qa_task)
    monkeypatch.setattr(qa_service.ai_service, "answer_question", answer_question)

    response = client.post(
        "/api/v1/qa/ask",
        json={"course_id": course["id"], "question": "我该如何复习《编译原理》？"},
        headers=student_headers,
    )

    assert response.status_code == 200, response.text
    data = response.json()["data"]
    assert calls == {"plan": 1, "answer": 1}
    assert data["is_out_of_scope"] is False
    assert "复习框架" in data["answer"]
    assert any(item.startswith("课程：编译原理") for item in captured["contexts"])
    assert any("章节结构：" in item or "页面内容：" in item or "资料：" in item for item in captured["contexts"])


def test_lesson_activity_skips_non_teaching_pages(client):
    course, chapter, _lesson_id, _teacher_headers, _student_headers = bootstrap_course_with_material(client)

    with db_session.SessionLocal() as db:
        db_course = db.get(Course, course["id"])
        assert db_course is not None
        material = CourseMaterial(
            course_id=db_course.id,
            chapter_id=chapter["id"],
            uploader_id=db_course.teacher_id,
            title="含介绍页课件",
            category=MaterialCategory.COURSEWARE.value,
            material_type=MaterialType.PDF.value,
            size_bytes=128,
            original_filename="intro.pdf",
            storage_path="tests/intro.pdf",
            extracted_text="课程介绍与线性变换定义",
            parse_status=ProcessStatus.READY.value,
            vector_status=ProcessStatus.READY.value,
        )
        db.add(material)
        db.flush()
        lesson = Lesson(
            course_id=db_course.id,
            chapter_id=chapter["id"],
            material_id=material.id,
            title="含介绍页课时",
            summary="测试",
            page_count=2,
            status=LessonStatus.PUBLISHED.value,
        )
        db.add(lesson)
        db.flush()
        intro_page = LessonPage(
            lesson_id=lesson.id,
            page_number=1,
            page_title="课程介绍",
            page_text="授课教师：张老师。邮箱：teacher@example.com。课程群二维码。",
            script_text="授课教师和联系方式介绍。",
            script_status=ProcessStatus.READY.value,
        )
        teaching_page = LessonPage(
            lesson_id=lesson.id,
            page_number=2,
            page_title="线性变换定义",
            page_text="线性变换定义：保持向量加法和数乘结构的映射。矩阵可以表示线性变换。",
            script_text="本页讲解线性变换定义、性质和矩阵表示方法。",
            script_status=ProcessStatus.READY.value,
        )
        db.add_all([intro_page, teaching_page])
        db.flush()
        intro_page_id = intro_page.id
        teaching_page_id = teaching_page.id
        ensure_lesson_pedagogy_artifacts(db, lesson=lesson, pages=[intro_page, teaching_page])
        db.commit()
        activities = page_activity_payload(db, lesson_page_ids=[intro_page_id, teaching_page_id])
        intro_activities = activities.get(intro_page_id, [])
        teaching_activities = activities.get(teaching_page_id, [])

    assert intro_activities == []
    assert teaching_activities


def test_qa_falls_back_to_database_chunks_when_vector_store_is_readonly(client, monkeypatch):
    course, _chapter, _lesson_id, _teacher_headers, student_headers = bootstrap_course_with_material(client)

    from app.services import knowledge as knowledge_service
    from app.services import qa as qa_service

    captured: dict[str, list[str]] = {}
    readonly_error = RuntimeError("Error updating collection: Database error: error returned from database: (code: 1032) attempt to write a readonly database")
    monkeypatch.setattr(knowledge_service.vector_store, "query_course", lambda *args, **kwargs: (_ for _ in ()).throw(readonly_error))
    monkeypatch.setattr(knowledge_service.vector_store, "upsert_chunks", lambda *args, **kwargs: (_ for _ in ()).throw(readonly_error))
    monkeypatch.setattr(
        qa_service.ai_service,
        "plan_qa_task",
        lambda **kwargs: {
            "scope": "specific",
            "question_type": "concept",
            "keywords": ["矩阵", "线性变换"],
            "search_phrases": ["矩阵 线性变换"],
            "retrieval_query": "矩阵 线性变换",
            "tools": ["search_courseware", "quote_source"],
        },
    )

    def answer_question(**kwargs):
        captured["contexts"] = list(kwargs["contexts"])
        return "已根据数据库知识切片回答。", False, None

    monkeypatch.setattr(qa_service.ai_service, "answer_question", answer_question)
    response = client.post(
        "/api/v1/qa/ask",
        json={"course_id": course["id"], "question": "矩阵可以表示什么？"},
        headers=student_headers,
    )
    assert response.status_code == 200, response.text
    data = response.json()["data"]
    assert data["is_out_of_scope"] is False
    assert captured["contexts"]
    assert data["sources"]


def test_search_course_knowledge_generalizes_numeric_examples(client, monkeypatch):
    course, _chapter, lesson_id, _teacher_headers, _student_headers = bootstrap_course_with_material(client)

    from app.services import knowledge as knowledge_service

    with db_session.SessionLocal() as db:
        page = db.scalar(select(LessonPage).where(LessonPage.lesson_id == lesson_id, LessonPage.page_number == 1))
        assert page is not None
        page.page_title = "加法步骤"
        page.page_text = "1+2 的解决步骤：先识别两个加数，再按照加法规则求和。"
        page.script_text = "这个例子说明两个数相加时可以先看操作数，再执行加法。"
        db.add(page)

        chunks = list(db.scalars(select(KnowledgeChunk).where(KnowledgeChunk.lesson_page_id == page.id)))
        assert chunks
        for chunk in chunks:
            chunk.title = "加法步骤"
            chunk.content = (
                "资料：算术示例\n"
                "页码：第1页\n"
                "页面标题：加法步骤\n"
                "页面内容：1+2 的解决步骤：先识别两个加数，再按照加法规则求和。"
            )
            chunk.tokens = ["加法", "步骤", "求和", "1+2"]
            db.add(chunk)
        db.commit()

        monkeypatch.setattr(knowledge_service.vector_store, "query_course", lambda *args, **kwargs: [])
        monkeypatch.setattr(knowledge_service.vector_store, "upsert_chunks", lambda *args, **kwargs: None)

        chunks = knowledge_service.search_course_knowledge(
            db,
            course_id=course["id"],
            query="8+9 怎么做？",
            lesson_id=lesson_id,
            limit=5,
        )

    assert chunks
    assert any("1+2" in chunk.content for chunk in chunks)


def test_qa_page_keyword_context_generalizes_numeric_examples(client):
    course, _chapter, lesson_id, _teacher_headers, _student_headers = bootstrap_course_with_material(client)

    from app.services import qa as qa_service

    with db_session.SessionLocal() as db:
        page = db.scalar(select(LessonPage).where(LessonPage.lesson_id == lesson_id, LessonPage.page_number == 1))
        assert page is not None
        page.page_title = "加法步骤"
        page.page_text = "1+2 的解决步骤：先识别两个加数，再按照加法规则求和。"
        page.script_text = None
        db.add(page)
        db.commit()

        contexts, sources = qa_service._page_keyword_context(
            db,
            course_id=course["id"],
            query="8+9 怎么做？",
            lesson_id=lesson_id,
            limit=3,
        )

    assert contexts
    assert any("1+2" in item for item in contexts)
    assert any(source.get("page_number") == 1 for source in sources)


def test_qa_retries_with_rewritten_query_when_original_question_is_too_concrete(client, monkeypatch):
    course, _chapter, lesson_id, _teacher_headers, student_headers = bootstrap_course_with_material(client)

    from app.services import qa as qa_service

    with db_session.SessionLocal() as db:
        pages = list(db.scalars(select(LessonPage).where(LessonPage.lesson_id == lesson_id).order_by(LessonPage.page_number)))
        assert len(pages) >= 2
        first_page = pages[0]
        second_page = pages[1]
        first_page.page_title = "题型模板"
        first_page.page_text = "同题型解题模板：先分析条件，再匹配方法模板，最后按步骤作答。"
        first_page.script_text = None
        db.add(first_page)
        db.commit()
        second_page_id = second_page.id

    captured: dict[str, list[str]] = {}
    monkeypatch.setattr(qa_service, "search_course_knowledge", lambda *args, **kwargs: [])
    monkeypatch.setattr(
        qa_service.ai_service,
        "rewrite_retrieval_query",
        lambda **kwargs: "同题型 解题模板 步骤 方法",
    )

    def answer_question(**kwargs):
        captured["contexts"] = list(kwargs["contexts"])
        return "先分析条件，再匹配方法模板。", False, None

    monkeypatch.setattr(qa_service.ai_service, "answer_question", answer_question)

    response = client.post(
        "/api/v1/qa/ask",
        json={"course_id": course["id"], "lesson_page_id": second_page_id, "question": "把样例里的人名和数字换掉后，这种题怎么做？"},
        headers=student_headers,
    )

    assert response.status_code == 200, response.text
    data = response.json()["data"]
    assert data["is_out_of_scope"] is False
    assert any("同题型解题模板" in item for item in captured["contexts"])


def test_tutoring_guidance_uses_course_retrieval_and_rewrite_retry(client, monkeypatch):
    course, _chapter, _lesson_id, _teacher_headers, student_headers = bootstrap_course_with_material(client)

    from app.services import tutoring as tutoring_service

    problem_resp = client.post(
        "/api/v1/tutoring/problems/text",
        json={"course_id": course["id"], "text": "把样例里的数字换掉后，这一类题怎么做？"},
        headers=student_headers,
    )
    assert problem_resp.status_code == 200, problem_resp.text
    problem = problem_resp.json()["data"]

    calls: list[str] = []

    def fake_search(db, *, course_id, query, chapter_id=None, lesson_id=None, lesson_page_id=None, limit=5):
        calls.append(query)
        if "检索重点" in query:
            return [
                KnowledgeChunk(
                    id=999,
                    course_id=course_id,
                    material_id=None,
                    lesson_page_id=None,
                    chapter_id=None,
                    title="同题型模板",
                    content="同题型解题模板：先分析条件，再匹配方法模板，最后按步骤作答。",
                    tokens=None,
                    embedding=None,
                    source_meta={},
                )
            ]
        return []

    captured: dict[str, object] = {}
    monkeypatch.setattr(tutoring_service, "search_course_knowledge", fake_search)
    monkeypatch.setattr(
        tutoring_service.ai_service,
        "rewrite_retrieval_query",
        lambda **kwargs: "同题型 方法 模板 步骤\n检索重点：题型 方法 步骤",
    )

    def fake_generate_problem_guidance(**kwargs):
        captured["contexts"] = list(kwargs.get("contexts") or [])
        return "## 分步思路\n1. 先分析条件。\n2. 再匹配方法模板。"

    monkeypatch.setattr(tutoring_service.ai_service, "generate_problem_guidance", fake_generate_problem_guidance)

    # #16：新题目仅 level 1 可解锁（顺序解锁），本用例只验证检索改写，用 level 1 即可。
    guidance_resp = client.get(
        f"/api/v1/tutoring/problems/{problem['id']}/guidance",
        params={"level": 1},
        headers=student_headers,
    )

    assert guidance_resp.status_code == 200, guidance_resp.text
    assert len(calls) >= 2
    assert any("检索重点" in query for query in calls)
    assert any("同题型解题模板" in item for item in captured["contexts"])


def test_tutoring_guidance_returns_existing_when_parallel_insert_wins(client, monkeypatch):
    course, _chapter, _lesson_id, _teacher_headers, student_headers = bootstrap_course_with_material(client)

    from app.services import tutoring as tutoring_service

    problem_resp = client.post(
        "/api/v1/tutoring/problems/text",
        json={"course_id": course["id"], "text": "根据语法制导定义计算表达式 8+7*12。"},
        headers=student_headers,
    )
    assert problem_resp.status_code == 200, problem_resp.text
    problem = problem_resp.json()["data"]
    inserted = {"value": False}

    monkeypatch.setattr(
        tutoring_service,
        "search_course_knowledge",
        lambda *args, **kwargs: [
            KnowledgeChunk(
                id=998,
                course_id=course["id"],
                material_id=None,
                lesson_page_id=None,
                chapter_id=None,
                title="语法制导定义",
                content="语法制导定义可以通过属性和语义规则计算表达式。",
                tokens=None,
                embedding=None,
                source_meta={},
            )
        ],
    )

    def generate_problem_guidance(**kwargs):
        if not inserted["value"]:
            with db_session.SessionLocal() as db:
                db.add(
                    ProblemGuidance(
                        problem_id=problem["id"],
                        level=1,
                        content="并发请求已经生成的辅导内容",
                        similar_questions=["已有相似题"],
                    )
                )
                db.commit()
            inserted["value"] = True
        return "较慢请求生成的辅导内容"

    monkeypatch.setattr(tutoring_service.ai_service, "generate_problem_guidance", generate_problem_guidance)

    # #16：新题目首个可解锁层级为 1，本用例验证并发插入命中唯一约束的回退，用 level 1。
    guidance_resp = client.get(
        f"/api/v1/tutoring/problems/{problem['id']}/guidance",
        params={"level": 1},
        headers=student_headers,
    )

    assert guidance_resp.status_code == 200, guidance_resp.text
    assert guidance_resp.json()["data"]["content"] == "并发请求已经生成的辅导内容"


def test_qa_can_answer_out_of_scope_with_notice_when_course_setting_enabled(client, monkeypatch):
    course, _chapter, _lesson_id, _teacher_headers, student_headers = bootstrap_course_with_material(client)

    from app.services import qa as qa_service

    with db_session.SessionLocal() as db:
        db_course = db.get(Course, course["id"])
        assert db_course is not None
        db_course.allow_general_ai_answer = True
        db.add(db_course)
        db.commit()

    monkeypatch.setattr(qa_service, "search_course_knowledge", lambda *args, **kwargs: [])
    monkeypatch.setattr(
        qa_service.ai_service,
        "plan_qa_task",
        lambda **kwargs: {
            "scope": "specific",
            "question_type": "specific",
            "keywords": ["背景知识"],
            "search_phrases": ["背景知识"],
            "retrieval_query": "背景知识",
            "tools": ["search_courseware", "quote_source"],
        },
    )
    monkeypatch.setattr(qa_service.ai_service, "rewrite_retrieval_query", lambda **kwargs: kwargs["question"])
    monkeypatch.setattr(
        qa_service.ai_service,
        "answer_general_question",
        lambda **kwargs: ("这是课程资料外的通用说明。", "通用推理"),
    )

    response = client.post(
        "/api/v1/qa/ask",
        json={"course_id": course["id"], "question": "这个课程没讲过的背景知识是什么？"},
        headers=student_headers,
    )

    assert response.status_code == 200, response.text
    data = response.json()["data"]
    assert data["is_out_of_scope"] is True
    assert "未在当前课程资料中检索到直接依据" in data["answer"]
    assert "通用说明" in data["answer"]


def test_tutoring_blocks_out_of_scope_guidance_when_course_setting_disabled(client, monkeypatch):
    course, _chapter, _lesson_id, _teacher_headers, student_headers = bootstrap_course_with_material(client)

    from app.services import tutoring as tutoring_service

    monkeypatch.setattr(tutoring_service, "search_course_knowledge", lambda *args, **kwargs: [])
    monkeypatch.setattr(tutoring_service.ai_service, "rewrite_retrieval_query", lambda **kwargs: kwargs["question"])

    guidance_called = {"value": False}

    def fake_generate_problem_guidance(**kwargs):
        guidance_called["value"] = True
        return "不应走到这里"

    monkeypatch.setattr(tutoring_service.ai_service, "generate_problem_guidance", fake_generate_problem_guidance)

    problem_resp = client.post(
        "/api/v1/tutoring/problems/text",
        json={"course_id": course["id"], "text": "这道明显不在资料里的题怎么做？"},
        headers=student_headers,
    )
    assert problem_resp.status_code == 200, problem_resp.text
    problem = problem_resp.json()["data"]

    guidance_resp = client.get(
        f"/api/v1/tutoring/problems/{problem['id']}/guidance",
        params={"level": 1},
        headers=student_headers,
    )

    assert guidance_resp.status_code == 200, guidance_resp.text
    assert guidance_called["value"] is False
    assert "未开启“资料外也可回答”" in guidance_resp.json()["data"]["content"]


def test_qa_image_attachment_uploads_and_participates_in_stream_answer(client, monkeypatch):
    course, _chapter, _lesson_id, _teacher_headers, student_headers = bootstrap_course_with_material(client)

    from app.services import qa as qa_service
    from app.services.ai import ChatDelta

    captured: dict[str, str] = {}
    monkeypatch.setattr(qa_service.ocr_service, "recognize", lambda upload, db=None: "矩阵可以表示线性变换")

    def stream_answer(**kwargs):
        captured["question"] = kwargs["question"]
        yield ChatDelta("content", "图片中提到矩阵可以表示线性变换。")

    monkeypatch.setattr(qa_service.ai_service, "stream_answer_question", stream_answer)

    upload_resp = client.post(
        "/api/v1/qa/attachments/image",
        data={"course_id": str(course["id"])},
        files={"file": ("matrix.png", PNG_BYTES, "image/png")},
        headers=student_headers,
    )
    assert upload_resp.status_code == 200, upload_resp.text
    attachment = upload_resp.json()["data"]
    assert attachment["filename"] == "matrix.png"
    assert attachment["ocr_text"] == "矩阵可以表示线性变换"

    with client.stream(
        "POST",
        "/api/v1/qa/ask/stream",
        json={"course_id": course["id"], "question": "帮我看图", "attachments": [attachment]},
        headers=student_headers,
    ) as response:
        body = "".join(response.iter_text())

    assert response.status_code == 200, body
    assert "event: error" not in body
    assert "图片中提到矩阵" in body
    assert "matrix.png" in body
    # OCR 文本以“不可信数据围栏”的形式进入 prompt（防提示注入），内容仍完整可用
    assert "矩阵可以表示线性变换" in captured["question"]
    assert "<<<IMAGE_OCR_START>>>" in captured["question"]
    assert "不可信" in captured["question"]
    assert "矩阵可以表示线性变换" in captured["question"]

    history_resp = client.get("/api/v1/qa/history", params={"course_id": course["id"]}, headers=student_headers)
    assert history_resp.status_code == 200, history_resp.text
    assert history_resp.json()["data"][0]["attachments"][0]["filename"] == "matrix.png"


def test_knowledge_points_use_local_explanations(client, monkeypatch):
    from app.services import knowledge as knowledge_service

    monkeypatch.setattr(knowledge_service.ai_service, "extract_knowledge_points", lambda text, db=None: ["矩阵"])

    def fail_explanation_call(**kwargs):
        raise AssertionError("knowledge explanations should not call the model synchronously")

    monkeypatch.setattr(knowledge_service.ai_service, "generate_knowledge_explanation", fail_explanation_call)
    course, chapter, _lesson_id, _teacher_headers, student_headers = bootstrap_course_with_material(client)

    response = client.get(
        "/api/v1/learning/knowledge-points",
        params={"course_id": course["id"], "chapter_id": chapter["id"]},
        headers=student_headers,
    )
    assert response.status_code == 200, response.text
    point = response.json()["data"][0]
    assert point["name"] == "矩阵"
    assert point["content_by_level"]["standard"]["definition"].startswith("矩阵")


def test_ebbinghaus_wrong_review_schedule():
    """艾宾浩斯遗忘曲线调度：连对按 [1,2,4,7,15,30] 天推进，走完整条曲线才判已掌握；到期判定 is_due。"""
    from datetime import UTC, datetime, timedelta

    from app.db.models import WrongQuestion
    from app.services.learning import (
        EBBINGHAUS_INTERVAL_DAYS,
        _apply_wrong_review_schedule,
        wrong_question_review_state,
    )

    now = datetime(2026, 1, 1, tzinfo=UTC)
    wrong = WrongQuestion(user_id=1, question_id=1, course_id=1, wrong_count=1, correct_streak=0)

    # 新错题/答错：排到次日（曲线第 0 档 = 1 天），未掌握。
    _apply_wrong_review_schedule(wrong, now)
    assert wrong.is_resolved is False
    assert wrong.next_review_at == now + timedelta(days=EBBINGHAUS_INTERVAL_DAYS[0])

    # 逐次答对：每档按曲线推进，走完前不归档。
    for streak in range(1, len(EBBINGHAUS_INTERVAL_DAYS)):
        wrong.correct_streak = streak
        _apply_wrong_review_schedule(wrong, now)
        assert wrong.is_resolved is False
        assert wrong.next_review_at == now + timedelta(days=EBBINGHAUS_INTERVAL_DAYS[streak])

    # 连对走完整条曲线：判已掌握、不再安排复习。
    wrong.correct_streak = len(EBBINGHAUS_INTERVAL_DAYS)
    _apply_wrong_review_schedule(wrong, now)
    assert wrong.is_resolved is True
    assert wrong.resolved_at == now
    assert wrong.next_review_at is None

    # is_due：到期(next_review <= now)判待复习；未来不判；已掌握不判。
    due = WrongQuestion(user_id=1, question_id=2, course_id=1, wrong_count=1, correct_streak=0)
    due.next_review_at = now - timedelta(days=1)
    assert wrong_question_review_state(due, now)["is_due"] is True
    due.next_review_at = now + timedelta(days=1)
    assert wrong_question_review_state(due, now)["is_due"] is False
    resolved = WrongQuestion(user_id=1, question_id=3, course_id=1, wrong_count=1, correct_streak=6, is_resolved=True)
    resolved.next_review_at = None
    state = wrong_question_review_state(resolved, now)
    assert state["is_due"] is False
    assert state["review_total"] == len(EBBINGHAUS_INTERVAL_DAYS)
