from io import BytesIO
from pathlib import Path
import sys
from types import SimpleNamespace

import fitz
from docx import Document
from pptx import Presentation
from sqlalchemy import select

from app.db import session as db_session
from app.db.models import AsyncTaskLog, CourseMaterial, KnowledgeChunk, Lesson, LessonPage
from app.services.parser import _localize_markdown_images, doc_parser_service, parse_material
from app.services.tts import markdown_to_speech_text, tts_service
from app.services.vector_store import vector_store


def register_user(client, *, email, password, nickname, role, student_no=None, employee_no=None):
    payload = {
        "email": email,
        "password": password,
        "nickname": nickname,
        "role": role,
        "student_no": student_no,
        "employee_no": employee_no,
    }
    if role == "teacher":
        admin_login = login_user(client, email="admin@classagent.com", password="Admin123456")
        response = client.post("/api/v1/admin/users/admin", json=payload, headers=auth_headers(admin_login["access_token"]))
    else:
        response = client.post("/api/v1/auth/register", json=payload)
    assert response.status_code == 200, response.text
    return response.json()["data"]


def login_user(client, *, email, password):
    response = client.post("/api/v1/auth/login", json={"email": email, "password": password})
    assert response.status_code == 200, response.text
    return response.json()["data"]


def auth_headers(token: str) -> dict[str, str]:
    return {"Authorization": f"Bearer {token}"}


def build_pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "极限定义"
    slide.placeholders[1].text = "极限描述函数在某点附近的变化趋势。"
    slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide2.shapes.title.text = "连续性"
    slide2.placeholders[1].text = "连续函数在区间内没有跳跃。"
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def build_docx_bytes() -> bytes:
    document = Document()
    document.add_heading("导数概念", level=1)
    document.add_paragraph("导数表示函数变化率。")
    document.add_paragraph("几何意义是切线斜率。")
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def build_pdf_bytes() -> bytes:
    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "积分用于度量累积量。")
    return pdf.tobytes()


def build_txt_bytes() -> bytes:
    return "矩阵可以表示线性变换。\n\n行列式反映缩放系数。".encode("utf-8")


def test_document_parsers_cover_all_supported_types(client, tmp_path: Path):
    pptx_path = tmp_path / "demo.pptx"
    pptx_path.write_bytes(build_pptx_bytes())
    pdf_path = tmp_path / "demo.pdf"
    pdf_path.write_bytes(build_pdf_bytes())
    docx_path = tmp_path / "demo.docx"
    docx_path.write_bytes(build_docx_bytes())
    txt_path = tmp_path / "demo.txt"
    txt_path.write_bytes(build_txt_bytes())

    with db_session.SessionLocal() as db:
        assert len(parse_material(pptx_path, "pptx", db=db, filename="demo.pptx")) == 2
        assert len(parse_material(pdf_path, "pdf", db=db, filename="demo.pdf")) == 1
        assert len(parse_material(docx_path, "docx", db=db, filename="demo.docx")) >= 1
        assert len(parse_material(txt_path, "txt", db=db, filename="demo.txt")) == 2


def test_doc_parser_layouts_group_into_ordered_pages():
    pages = doc_parser_service._pages_from_layouts(
        [
            {"pageNum": 1, "index": 2, "markdownContent": "第二页正文"},
            {"pageNum": 0, "index": 1, "type": "text", "text": "第一页正文"},
            {"pageNum": 0, "index": 0, "type": "title", "text": "第一页标题", "markdownContent": "# 第一页标题"},
        ]
    )
    assert len(pages) == 2
    assert pages[0]["page_number"] == 1
    assert pages[0]["page_title"] == "第一页标题"
    assert pages[0]["page_text"].startswith("# 第一页标题")
    assert pages[1]["page_number"] == 2
    assert "第二页正文" in pages[1]["page_text"]


def test_docmind_markdown_images_are_persisted_to_configured_storage(monkeypatch):
    calls: dict[str, object] = {}

    class FakeResponse:
        headers = {"content-type": "image/png"}
        content = b"fake-png"

        def raise_for_status(self) -> None:
            return None

    class FakeClient:
        def __init__(self, *args, **kwargs):
            return None

        def __enter__(self):
            return self

        def __exit__(self, *args):
            return False

        def get(self, url):
            calls["download_url"] = url
            return FakeResponse()

    monkeypatch.setattr("app.services.parser.httpx.Client", FakeClient)
    monkeypatch.setattr(
        "app.services.parser.storage_service.save_bytes",
        lambda content, *, folder, filename, db=None: (
            calls.update({"content": content, "folder": folder, "filename": filename})
            or f"{folder}/{filename}"
        ),
    )
    monkeypatch.setattr(
        "app.services.parser.storage_service.public_url",
        lambda relative_path, db=None: f"https://cdn.example.com/{relative_path}",
    )

    temporary_url = (
        "http://docmind-api-cn-hangzhou.oss-cn-hangzhou.aliyuncs.com/publicDocStreamStructure/10.png"
        "?Expires=1777819773&OSSAccessKeyId=STS.demo&Signature=abc%3D&security-token=token"
    )
    content = f"第一页\n![公式截图]({temporary_url})"

    rewritten = _localize_markdown_images(content, db=None, cache={})

    assert "docmind-api-cn-hangzhou" not in rewritten
    assert "OSSAccessKeyId" not in rewritten
    assert "https://cdn.example.com/docmind_images/" in rewritten
    assert calls["download_url"] == temporary_url
    assert calls["content"] == b"fake-png"
    assert str(calls["folder"]).startswith("docmind_images/")


def test_expired_docmind_markdown_images_are_not_kept(monkeypatch):
    class FakeClient:
        def __init__(self, *args, **kwargs):
            return None

        def __enter__(self):
            return self

        def __exit__(self, *args):
            return False

        def get(self, url):
            raise RuntimeError("must not use non-http errors")

    class FailingResponse:
        headers = {}
        content = b""

        def raise_for_status(self):
            import httpx

            raise httpx.HTTPStatusError("forbidden", request=None, response=None)

    class ExpiredClient(FakeClient):
        def get(self, url):
            return FailingResponse()

    monkeypatch.setattr("app.services.parser.httpx.Client", ExpiredClient)
    temporary_url = (
        "https://docmind-api-cn-hangzhou.oss-cn-hangzhou.aliyuncs.com/publicDocStreamStructure/10.png"
        "?Expires=1&OSSAccessKeyId=STS.demo&Signature=abc"
    )

    rewritten = _localize_markdown_images(f"说明 ![图1]({temporary_url}) 结束", db=None, cache={})

    assert temporary_url not in rewritten
    assert "OSSAccessKeyId" not in rewritten
    assert "图1" in rewritten


def test_markdown_to_speech_text_removes_markdown_and_links():
    cleaned = markdown_to_speech_text(
        """
        # 第一章
        - **重点**：[语法分析](https://example.com/a)
        ![流程图](https://example.com/image.png)
        | A | B |
        | --- | --- |
        | `归约` | $\\frac{a}{b}$ |
        > \\mathbb{P}_{2}=\\left\\{\\begin{array}{c c}{\\Gamma_B}\\end{array}\\right.
        """
    )

    assert "第一章" in cleaned
    assert "语法分析" in cleaned
    assert "流程图" in cleaned
    assert "https://" not in cleaned
    for marker in ("#", "**", "[", "](", "![", "|", "`", "\\mathbb", "\\frac"):
        assert marker not in cleaned


def test_aliyun_tts_uses_official_nls_sdk(monkeypatch, tmp_path: Path):
    calls: dict[str, object] = {}
    tts_service._token_cache.clear()

    class FakeSynthesizer:
        def __init__(self, **kwargs):
            calls["init"] = kwargs
            self.on_data = kwargs["on_data"]

        def start(self, **kwargs):
            calls["start"] = kwargs
            self.on_data(b"fake-audio")

    monkeypatch.setitem(sys.modules, "nls", SimpleNamespace(NlsSpeechSynthesizer=FakeSynthesizer))
    monkeypatch.setattr(
        tts_service,
        "_create_token",
        lambda config: (calls.update({"token_config": config}) or ("generated-token", 4_102_444_800)),
    )
    monkeypatch.setattr(
        "app.services.tts.storage_service.save_bytes",
        lambda content, *, folder, filename, db=None: (calls.update({"content": content, "filename": filename}) or "generated/audio/fake.wav"),
    )
    monkeypatch.setattr("app.services.tts.storage_service.public_url", lambda relative_path, db=None: f"/static/{relative_path}")

    url, duration = tts_service._synthesize_aliyun(
        "连接测试",
        None,
        {"access_key_id": "ak", "access_key_secret": "secret", "appkey": "app", "voice": "xiaoyun", "format": "wav"},
    )

    assert calls["token_config"]["access_key_id"] == "ak"
    assert calls["init"]["appkey"] == "app"
    assert calls["init"]["token"] == "generated-token"
    assert calls["init"]["url"] == "wss://nls-gateway.cn-shanghai.aliyuncs.com/ws/v1"
    assert calls["start"]["voice"] == "xiaoyun"
    assert calls["content"] == b"fake-audio"
    assert url == "/static/generated/audio/fake.wav"
    assert duration >= 2


def test_material_management_flow(client):
    register_user(
        client,
        email="teacher2@example.com",
        password="Teacher123",
        nickname="王老师",
        role="teacher",
        employee_no="T2026002",
    )
    register_user(
        client,
        email="student2@example.com",
        password="Student123",
        nickname="赵同学",
        role="student",
        student_no="S2026002",
    )
    teacher_login = login_user(client, email="teacher2@example.com", password="Teacher123")
    teacher_headers = auth_headers(teacher_login["access_token"])
    student_login = login_user(client, email="student2@example.com", password="Student123")
    student_headers = auth_headers(student_login["access_token"])

    course_resp = client.post(
        "/api/v1/courses",
        json={"name": "线性代数", "description": "矩阵与向量", "term": "2026春"},
        headers=teacher_headers,
    )
    assert course_resp.status_code == 200, course_resp.text
    course = course_resp.json()["data"]

    chapter_resp = client.post(
        f"/api/v1/courses/{course['id']}/chapters",
        json={"title": "第一章 矩阵", "description": "基础内容", "order_index": 1},
        headers=teacher_headers,
    )
    assert chapter_resp.status_code == 200, chapter_resp.text
    chapter = chapter_resp.json()["data"]

    join_resp = client.post(
        "/api/v1/courses/join",
        json={"course_code": course["course_code"]},
        headers=student_headers,
    )
    assert join_resp.status_code == 200, join_resp.text

    upload_resp = client.post(
        "/api/v1/materials",
        data={
            "course_id": str(course["id"]),
            "title": "矩阵基础课件",
            "category": "courseware",
            "chapter_id": str(chapter["id"]),
        },
        files={
            "file": (
                "matrix.pptx",
                build_pptx_bytes(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
        headers=teacher_headers,
    )
    assert upload_resp.status_code == 200, upload_resp.text
    material = upload_resp.json()["data"]
    assert material["parse_status"] == "ready"
    assert material["vector_status"] == "ready"
    assert material["preview_url"].startswith("/static/")

    detail_resp = client.get(f"/api/v1/materials/{material['id']}", headers=teacher_headers)
    assert detail_resp.status_code == 200, detail_resp.text
    detail = detail_resp.json()["data"]
    assert detail["material"]["preview_url"].startswith("/static/")
    assert detail["lesson_page_count"] == 2
    assert len(detail["pages"]) == 2
    assert detail["pages"][0]["script_text"]
    assert detail["pages"][0]["audio_url"].startswith("/static/")
    assert detail["pages"][0]["audio_url"].endswith(".wav")
    with db_session.SessionLocal() as db:
        stored_material = db.get(CourseMaterial, material["id"])
        stored_page = db.get(LessonPage, detail["pages"][0]["id"])
        stored_material.preview_url = "http://127.0.0.1:8000/static/uploads/legacy/demo.pptx"
        stored_page.audio_url = "http://127.0.0.1:8000/static/generated/audio/legacy.wav"
        db.add_all([stored_material, stored_page])
        db.commit()
        chunks = list(db.scalars(select(KnowledgeChunk).where(KnowledgeChunk.material_id == material["id"])))
        assert len(chunks) == 2
        assert all(isinstance(chunk.embedding, list) and chunk.embedding for chunk in chunks)
        vector_rows = vector_store.query_course(db, course_id=course["id"], query="函数变化趋势怎样理解", limit=2)
        assert vector_rows

    legacy_detail_resp = client.get(f"/api/v1/materials/{material['id']}", headers=teacher_headers)
    assert legacy_detail_resp.status_code == 200, legacy_detail_resp.text
    legacy_detail = legacy_detail_resp.json()["data"]
    assert legacy_detail["material"]["preview_url"] == "/static/uploads/legacy/demo.pptx"
    assert legacy_detail["pages"][0]["audio_url"] == "/static/generated/audio/legacy.wav"

    page_id = detail["pages"][0]["id"]
    update_script_resp = client.patch(
        f"/api/v1/materials/pages/{page_id}/script",
        json={"script_text": "这是教师手动修订后的讲解脚本。"},
        headers=teacher_headers,
    )
    assert update_script_resp.status_code == 200, update_script_resp.text
    assert update_script_resp.json()["data"]["script_text"] == "这是教师手动修订后的讲解脚本。"

    regenerate_resp = client.post(
        f"/api/v1/materials/pages/{page_id}/script/regenerate",
        headers=teacher_headers,
    )
    assert regenerate_resp.status_code == 200, regenerate_resp.text
    assert regenerate_resp.json()["data"]["audio_url"].endswith(".wav")

    student_list_resp = client.get(
        "/api/v1/materials",
        params={"course_id": course["id"], "keyword": "矩阵"},
        headers=student_headers,
    )
    assert student_list_resp.status_code == 200, student_list_resp.text
    assert len(student_list_resp.json()["data"]) == 1

    reprocess_resp = client.post(
        f"/api/v1/materials/{material['id']}/reprocess",
        headers=teacher_headers,
    )
    assert reprocess_resp.status_code == 200, reprocess_resp.text

    update_material_resp = client.patch(
        f"/api/v1/materials/{material['id']}",
        json={"title": "矩阵基础课件-修订版", "category": "handout", "chapter_id": chapter["id"]},
        headers=teacher_headers,
    )
    assert update_material_resp.status_code == 200, update_material_resp.text
    assert update_material_resp.json()["data"]["title"] == "矩阵基础课件-修订版"

    delete_resp = client.delete(f"/api/v1/materials/{material['id']}", headers=teacher_headers)
    assert delete_resp.status_code == 200, delete_resp.text


def test_material_pipeline_keeps_parse_ready_when_tts_fails(client, monkeypatch):
    register_user(
        client,
        email="teacher-tts-fallback@example.com",
        password="Teacher123",
        nickname="语音降级老师",
        role="teacher",
        employee_no="TTS-FALLBACK",
    )
    teacher_login = login_user(client, email="teacher-tts-fallback@example.com", password="Teacher123")
    teacher_headers = auth_headers(teacher_login["access_token"])
    course_resp = client.post(
        "/api/v1/courses",
        json={"name": "语音降级课程", "description": "测试", "term": "2026春"},
        headers=teacher_headers,
    )
    assert course_resp.status_code == 200, course_resp.text
    course = course_resp.json()["data"]

    monkeypatch.setattr(tts_service, "synthesize", lambda *args, **kwargs: (_ for _ in ()).throw(RuntimeError("tts down")))

    upload_resp = client.post(
        "/api/v1/materials",
        data={"course_id": str(course["id"]), "title": "语音失败课件", "category": "courseware"},
        files={"file": ("fallback.txt", build_txt_bytes(), "text/plain")},
        headers=teacher_headers,
    )
    assert upload_resp.status_code == 200, upload_resp.text
    material = upload_resp.json()["data"]
    assert material["parse_status"] == "ready"

    detail_resp = client.get(f"/api/v1/materials/{material['id']}", headers=teacher_headers)
    assert detail_resp.status_code == 200, detail_resp.text
    detail = detail_resp.json()["data"]
    assert detail["lesson_page_count"] == 2
    assert detail["pages"][0]["script_text"]
    assert detail["pages"][0]["audio_url"] is None
    with db_session.SessionLocal() as db:
        task = db.scalar(select(AsyncTaskLog).where(AsyncTaskLog.target_id == material["id"]).order_by(AsyncTaskLog.id.desc()))
        assert task.status == "ready"
        assert "语音合成失败" in "；".join(task.detail["warnings"])


def test_material_pipeline_keeps_pages_when_vector_store_fails(client, monkeypatch):
    register_user(
        client,
        email="teacher-vector-fallback@example.com",
        password="Teacher123",
        nickname="向量降级老师",
        role="teacher",
        employee_no="VECTOR-FALLBACK",
    )
    teacher_login = login_user(client, email="teacher-vector-fallback@example.com", password="Teacher123")
    teacher_headers = auth_headers(teacher_login["access_token"])
    course_resp = client.post(
        "/api/v1/courses",
        json={"name": "向量降级课程", "description": "测试", "term": "2026春"},
        headers=teacher_headers,
    )
    assert course_resp.status_code == 200, course_resp.text
    course = course_resp.json()["data"]

    monkeypatch.setattr(vector_store, "upsert_chunks", lambda *args, **kwargs: (_ for _ in ()).throw(RuntimeError("readonly vector db")))

    upload_resp = client.post(
        "/api/v1/materials",
        data={"course_id": str(course["id"]), "title": "向量失败课件", "category": "courseware"},
        files={"file": ("vector.txt", build_txt_bytes(), "text/plain")},
        headers=teacher_headers,
    )
    assert upload_resp.status_code == 200, upload_resp.text
    material = upload_resp.json()["data"]
    assert material["parse_status"] == "ready"
    assert material["vector_status"] == "failed"

    detail_resp = client.get(f"/api/v1/materials/{material['id']}", headers=teacher_headers)
    assert detail_resp.status_code == 200, detail_resp.text
    detail = detail_resp.json()["data"]
    assert detail["lesson_page_count"] == 2
    assert detail["pages"][0]["script_text"]
    assert detail["pages"][0]["audio_url"].startswith("/static/")


def test_failed_material_with_existing_pages_is_recovered_for_teacher_views(client):
    register_user(
        client,
        email="teacher-repair@example.com",
        password="Teacher123",
        nickname="恢复老师",
        role="teacher",
        employee_no="REPAIR-MATERIAL",
    )
    teacher_login = login_user(client, email="teacher-repair@example.com", password="Teacher123")
    teacher_headers = auth_headers(teacher_login["access_token"])
    course_resp = client.post(
        "/api/v1/courses",
        json={"name": "资料恢复课程", "description": "测试", "term": "2026春"},
        headers=teacher_headers,
    )
    assert course_resp.status_code == 200, course_resp.text
    course = course_resp.json()["data"]

    upload_resp = client.post(
        "/api/v1/materials",
        data={"course_id": str(course["id"]), "title": "旧状态失败课件", "category": "courseware"},
        files={"file": ("repair.txt", build_txt_bytes(), "text/plain")},
        headers=teacher_headers,
    )
    assert upload_resp.status_code == 200, upload_resp.text
    material = upload_resp.json()["data"]

    with db_session.SessionLocal() as db:
        stored_material = db.get(CourseMaterial, material["id"])
        lesson = db.scalar(select(Lesson).where(Lesson.material_id == material["id"]))
        assert stored_material is not None
        assert lesson is not None
        stored_material.parse_status = "failed"
        stored_material.vector_status = "failed"
        stored_material.extracted_text = None
        lesson.page_count = 0
        lesson.status = "draft"
        db.add_all([stored_material, lesson])
        db.commit()

    summary_resp = client.get(f"/api/v1/teacher/courses/{course['id']}/materials/summary", headers=teacher_headers)
    assert summary_resp.status_code == 200, summary_resp.text
    summary = summary_resp.json()["data"]
    assert summary["ready"] == 1
    assert summary["stats"]["by_status"]["ready"] == 1

    list_resp = client.get("/api/v1/materials", params={"course_id": course["id"]}, headers=teacher_headers)
    assert list_resp.status_code == 200, list_resp.text
    listed_material = list_resp.json()["data"][0]
    assert listed_material["parse_status"] == "ready"
    assert listed_material["vector_status"] == "failed"
    assert listed_material["extracted_text"]

    detail_resp = client.get(f"/api/v1/materials/{material['id']}", headers=teacher_headers)
    assert detail_resp.status_code == 200, detail_resp.text
    detail = detail_resp.json()["data"]
    assert detail["material"]["parse_status"] == "ready"
    assert detail["lesson_page_count"] == 2


def test_script_save_survives_tts_failure(client, monkeypatch):
    register_user(
        client,
        email="teacher-script-tts@example.com",
        password="Teacher123",
        nickname="脚本老师",
        role="teacher",
        employee_no="SCRIPT-TTS",
    )
    teacher_login = login_user(client, email="teacher-script-tts@example.com", password="Teacher123")
    teacher_headers = auth_headers(teacher_login["access_token"])
    course_resp = client.post(
        "/api/v1/courses",
        json={"name": "脚本保存课程", "description": "测试", "term": "2026春"},
        headers=teacher_headers,
    )
    assert course_resp.status_code == 200, course_resp.text
    course = course_resp.json()["data"]

    upload_resp = client.post(
        "/api/v1/materials",
        data={"course_id": str(course["id"]), "title": "脚本保存课件", "category": "courseware"},
        files={"file": ("script.txt", build_txt_bytes(), "text/plain")},
        headers=teacher_headers,
    )
    assert upload_resp.status_code == 200, upload_resp.text
    material = upload_resp.json()["data"]
    detail_resp = client.get(f"/api/v1/materials/{material['id']}", headers=teacher_headers)
    page_id = detail_resp.json()["data"]["pages"][0]["id"]

    monkeypatch.setattr(tts_service, "synthesize", lambda *args, **kwargs: (_ for _ in ()).throw(RuntimeError("tts down")))

    update_resp = client.patch(
        f"/api/v1/materials/pages/{page_id}/script",
        json={"script_text": "TTS 失败时也应该保存这段脚本。"},
        headers=teacher_headers,
    )
    assert update_resp.status_code == 200, update_resp.text
    page = update_resp.json()["data"]
    assert page["script_text"] == "TTS 失败时也应该保存这段脚本。"
    assert page["script_status"] == "ready"
    assert page["audio_url"] is None
